"""
Generate the Phase I — Review 2 (of 3) PPT — Unified Telecom Analyzer.

Strictly structured around the "Second Review" checklist for Phase I:
  * Title & Abstract
  * Detailed Algo Design
  * Contribution of the candidate
  * Results obtained (intermediate)
  * References
  * 80% of code Implementation (Phase 1)

Uses a distinct navy/gold/teal colour theme (different from the other decks).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ── Navy / Gold / Teal theme (distinct from cyan-on-dark-blue decks) ──
BG       = RGBColor(0x1B, 0x1A, 0x2E)   # deep indigo-navy
ACCENT   = RGBColor(0xFF, 0xC1, 0x07)   # gold
WHITE    = RGBColor(0xF1, 0xF1, 0xF6)
GREEN    = RGBColor(0x4A, 0xDE, 0x80)
TEAL     = RGBColor(0x2E, 0xC4, 0xB6)
CORAL    = RGBColor(0xE9, 0x45, 0x60)
GREY     = RGBColor(0xAE, 0xAE, 0xC2)
LILAC    = RGBColor(0x9B, 0x8C, 0xF2)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size, color, line_h=0.32):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_title(slide, text):
    textbox(slide, 0.3, 0.12, 9.4, 0.5, text, 22, ACCENT, bold=True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    return slide


# ── Slide 1: Title ────────────────────────────────────────────────────

def s_title(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 0.9, 9.0, 0.9,
            "Unified Telecom Analyzer", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.95, 9.0, 0.5,
            "Multi-Modal Anomaly Detection & Explanation Framework for 5G Networks",
            18, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.5, 9.0, 0.4,
            "PHASE I — REVIEW 2  (of 3)  |  Research Objective 1",
            16, TEAL, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.1, 9.0, 0.4,
            "M.Tech Data Science  |  PES University, Bangalore — Great Learning",
            13, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.55, 9.0, 0.4,
            "Candidate: Vikas  |  Repository: telecom-analyzer",
            12, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.6, 9.0, 0.4,
            "Python · pyshark · scikit-learn · PyTorch · Prophet · FAISS · Ollama · Streamlit · FastAPI",
            11, GREY, align=PP_ALIGN.CENTER)


# ── Slide 2: Title & Abstract ────────────────────────────────────────

def s_abstract(prs):
    slide = new_slide(prs)
    slide_title(slide, "Title & Abstract")
    textbox(slide, 0.3, 0.7, 9.4, 0.4,
            "Unified Telecom Analyzer: A Multi-Modal, Multi-Method Anomaly "
            "Detection & Explanation Framework for 5G Networks",
            14, TEAL, bold=True, wrap=True)
    text = (
        "5G core and radio-access networks generate enormous volumes of "
        "heterogeneous signalling and performance data — PCAP traces "
        "(NAS/NGAP/RRC/F1AP/E1AP/XnAP), DU/CU layer-1/2 statistics, and "
        "cell-level KPI time-series. Operators today inspect these sources "
        "in isolation with single-method tools, missing anomalies that only "
        "become visible when signalling, radio-stack counters and KPIs are "
        "viewed together.\n\n"
        "This project builds a single pipeline that parses all three data "
        "types against their respective 3GPP specifications, runs a "
        "complementary 18-detector ensemble (6 unsupervised/statistical "
        "methods x 3 domains) over them, correlates anomalies across "
        "sources on (cell, time), forecasts anomalies up to 4 hours ahead, "
        "and explains the most likely root cause using a local, "
        "3GPP-grounded RAG + LLM — with a deterministic rule-based "
        "fallback for offline use."
    )
    textbox(slide, 0.3, 1.25, 9.4, 4.0, text, 12, WHITE, wrap=True)


# ── Slide 3: Detailed Algorithm Design — Architecture ─────────────────

def s_algo_architecture(prs):
    slide = new_slide(prs)
    slide_title(slide, "Detailed Algo Design — 5-Layer Architecture")
    layers = [
        ("INPUT",
         "PCAP (.pcap/.pcapng)  ·  DU/CU Stats (srsRAN/OAI/NIST CSV/Parquet)  ·  KPI Time-series (Excel/CSV)",
         TEAL, 0.75),
        ("PARSING — TS 24.501 / 38.413 / 38.331 / 38.473 / 38.463 / 38.423",
         "NAS · NGAP · RRC · F1AP · E1AP · XnAP · KPI Parser · Stats Parser (srsRAN/OAI/NIST)",
         GREEN, 1.55),
        ("DETECTION — 18 detectors total",
         "PCAP: IF · Statistical · OC-SVM · LOF · Elliptic Env. · LSTM-AE   "
         "KPI: Threshold · Peer Comp. · Trend · IQR · CUSUM · Bollinger   "
         "Stats: same 6 methods on L1/L2 counters",
         ACCENT, 2.45),
        ("CORRELATION + PREDICTION + RAG/LLM + FEEDBACK",
         "Event Router (cross-source correlation)  |  Prophet + LSTM 4h-ahead forecast  |  "
         "FAISS+MiniLM RAG -> Ollama phi3:mini explanation  |  Feedback store -> nightly retrainer",
         LILAC, 3.35),
        ("OUTPUT",
         "Streamlit Dashboard (http://localhost:8501)  ·  FastAPI REST API (:8000)  ·  CLI Pipeline Runner",
         CORAL, 4.2),
    ]
    for label, items, color, top in layers:
        textbox(slide, 0.3, top, 2.7, 0.3, label, 9, color, bold=True)
        textbox(slide, 0.3, top + 0.28, 9.2, 0.5, items, 10, WHITE, wrap=True)


# ── Slide 4: Detailed Algorithm Design — PCAP Ensemble ─────────────────

def s_algo_pcap(prs):
    slide = new_slide(prs)
    slide_title(slide, "Detailed Algo Design — PCAP Anomaly Detection (6 methods)")
    headers = ["Method", "What it catches", "Why chosen"]
    rows = [
        ("Isolation Forest", "Global statistical outliers in procedure-level features",
         "Fast, tree-based, no distribution assumption — good first pass"),
        ("Statistical / Cascade (rule-based)", "Known SLA breaches & cascading failures",
         "Encodes telecom domain knowledge / 3GPP timer-based rules directly"),
        ("One-Class SVM", "Non-linear normal-region boundary violations",
         "Captures complex normal regions IF's axis-aligned splits miss"),
        ("Local Outlier Factor (LOF)", "Local density anomalies (peer-relative)",
         "Detects anomalies that are normal globally but odd vs. local cluster"),
        ("Elliptic Envelope", "Mahalanobis-distance outliers vs. Gaussian baseline",
         "Cheap, interpretable, good when features are roughly Gaussian"),
        ("LSTM Autoencoder", "Procedure-order / sequence violations",
         "Only sequence-aware method — catches out-of-order signalling"),
    ]
    col_x = [0.3, 2.6, 5.4]
    col_w = [2.2, 2.7, 4.0]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.75, col_w[i], 0.3, h, 11, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.12 + r * 0.72
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.68, val, 9.5, WHITE, wrap=True)
    textbox(slide, 0.3, 5.2, 9.2, 0.3,
            "Methods are complementary, not redundant — ensemble agreement raises confidence and lowers false positives.",
            10, TEAL, italic=True, wrap=True)


# ── Slide 5: Detailed Algorithm Design — KPI/Stats Ensemble ────────────

def s_algo_kpi(prs):
    slide = new_slide(prs)
    slide_title(slide, "Detailed Algo Design — KPI / L1-L2 Stats Detection (6 methods)")
    headers = ["Method", "What it catches", "Why chosen"]
    rows = [
        ("Threshold Violation", "Per-row breach of 3GPP-catalogue warning/critical limits",
         "Direct, explainable mapping to operator SLAs"),
        ("Peer Comparison (z-score)", "Cells under-performing vs. fleet average",
         "Surfaces relative degradation even within 'normal' absolute range"),
        ("Trend (Linear Regression)", "Monotonic degradation, slope > threshold/hr",
         "Catches slow drift before it crosses a hard threshold"),
        ("IQR (Tukey Fence)", "Distribution-free outliers, robust to skew",
         "No Gaussian assumption — works for skewed KPIs like latency"),
        ("CUSUM", "Cumulative-sum change-point detection",
         "Detects gradual drift 5+ hours before threshold breach"),
        ("Bollinger Bands", "Burst spikes vs. rolling envelope",
         "Lagged baseline catches sudden bursts that CUSUM smooths over"),
    ]
    col_x = [0.3, 2.6, 5.4]
    col_w = [2.2, 2.7, 4.0]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.75, col_w[i], 0.3, h, 11, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.12 + r * 0.72
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.68, val, 9.5, WHITE, wrap=True)
    textbox(slide, 0.3, 5.2, 9.2, 0.3,
            "Same 6-method ensemble re-applied to DU/CU L1/L2 stats (srsRAN/OAI/NIST) for protocol-level corroboration.",
            10, TEAL, italic=True, wrap=True)


# ── Slide 6: Detailed Algorithm Design — Correlation/Prediction/RAG/MLOps ──

def s_algo_pipeline(prs):
    slide = new_slide(prs)
    slide_title(slide, "Detailed Algo Design — Correlation, Prediction, RAG/LLM, MLOps")
    bullets(slide, 0.4, 0.75, 9.2, 4.6, [
        "Event Router: groups anomalies from PCAP, KPI and Stats by "
        "(cell, category) within a 1-hour window; >=2 sources agreeing on "
        "the same cell/category is marked 'cross-source correlated'.",
        "Prediction Layer: Prophet (seasonality/trend decomposition) + LSTM "
        "(sequence model) forecast KPI degradation up to 4 hours ahead; "
        "predicted anomalies share the same schema as live detections "
        "(state='predicted', lead_time_h).",
        "RAG Retrieval: 38 chunks from 6 3GPP specs (TS 24.501, 38.413, "
        "38.331, 38.473, 38.463, 38.423) embedded with all-MiniLM-L6-v2, "
        "indexed in FAISS for nearest-neighbour spec lookup per anomaly.",
        "LLM Explanation: Ollama phi3:mini generates a root-cause "
        "explanation grounded in the retrieved spec text; deterministic "
        "rule-based template is used as fallback when the LLM is "
        "unavailable (offline/no-GPU).",
        "MLOps Feedback Loop: engineers mark detections as true/false "
        "positive -> feedback store (JSONL) -> nightly retrainer adjusts "
        "per-detector thresholds/contamination using the observed "
        "false-positive rate (proportional controller, not full retraining).",
    ], 12, WHITE, line_h=0.5)


# ── Slide 7: Contribution of the Candidate ────────────────────────────

def s_contribution(prs):
    slide = new_slide(prs)
    slide_title(slide, "Contribution of the Candidate")
    bullets(slide, 0.4, 0.78, 9.2, 4.4, [
        "Designed and implemented the full 3GPP-aligned parsing layer "
        "(PCAP via pyshark + raw-discriminator fallback; DU/CU stats for "
        "srsRAN/OAI/NIST/Parquet; KPI Excel/CSV vs. a hand-built 3GPP KPI "
        "catalogue with thresholds/directions/units).",
        "Built and validated the 18-detector ensemble (6 methods x 3 "
        "domains) from scratch, each returning severity, rationale, and "
        "evidence — plus a cross-method comparison matrix justifying why "
        "each method is needed.",
        "Designed the cross-source Event Router (correlation key, "
        "correlation window, severity aggregation) and the Prophet+LSTM "
        "4-hour prediction layer sharing a common anomaly schema.",
        "Built the local RAG pipeline end-to-end (spec chunking, MiniLM "
        "embeddings, FAISS index) and wired it to Ollama phi3:mini with a "
        "deterministic fallback explainer.",
        "Implemented the MLOps feedback loop (feedback store + nightly "
        "retrainer as a proportional controller on detector thresholds).",
        "Exposed everything via a Streamlit dashboard, FastAPI REST API, "
        "and CLI orchestrator, and wrote 56 automated tests (all passing) "
        "covering parsers, detectors, router, prediction and retrainer.",
    ], 11.5, WHITE, line_h=0.48)


# ── Slide 8: Results Obtained (Intermediate) ───────────────────────────

def s_results(prs):
    slide = new_slide(prs)
    slide_title(slide, "Results Obtained (Intermediate)")
    bullets(slide, 0.4, 0.7, 9.2, 1.0, [
        "Synthetic data covering known-good + injected-failure scenarios "
        "per data type; figures below are from actual end-to-end pipeline "
        "runs (--no-llm) on these datasets.",
    ], 11, WHITE, line_h=0.35)

    headers = ["Dataset", "Size", "Anomalies (ensemble)", "Key findings"]
    rows = [
        ("ue_attach_10.pcap", "199 pkts / 14 procedures", "12",
         "NAS_Registration & XnAP_HandoverRequest flagged by 4-6 of 6 detectors (IF, OC-SVM, LOF, LSTM-AE)"),
        ("ue_handover_20.pcap", "108 pkts / 6 procedures", "6",
         "NGAP_HandoverPreparation, XnAP_HandoverRequest, RRC_Reconfiguration flagged by IF/OC-SVM/LSTM-AE"),
        ("kpi_10hr_sample.csv", "6000 rows x 25 KPIs", "1486",
         "Threshold 1057, Bollinger 139, CUSUM 176, IQR 97, Peer 17 — congestion/outage/drift injected & recovered"),
        ("srsran_stats.csv", "960 rows x 15 cols", "29",
         "Trend detector (29) — slow drift captured in L1/L2 counters; other 5 methods clean"),
    ]
    col_x = [0.3, 2.3, 4.2, 5.9]
    col_w = [1.9, 1.8, 1.6, 3.3]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 1.55, col_w[i], 0.3, h, 11, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.92 + r * 0.62
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.6, val, 9.5, WHITE, wrap=True)

    textbox(slide, 0.3, 4.55, 9.2, 0.3,
            "Performance (--no-llm, CPU-only, 16GB WSL2): PCAP attach 6.4s, "
            "PCAP handover 1.4s, KPI (6000x25) 11.6s, Stats (960x15) 1.1s.",
            10, GREY, wrap=True)
    textbox(slide, 0.3, 4.95, 9.2, 0.3,
            "Predicted anomalies (Prophet/LSTM) verified to use the same "
            "output schema as live detections (state='predicted', lead_time_h).",
            10, GREY, wrap=True)


# ── Slide 9: 80% of Code Implementation (Phase 1) ─────────────────────

def s_code_status(prs):
    slide = new_slide(prs)
    slide_title(slide, "80% of Code Implementation (Phase 1)")
    headers = ["Component", "Status"]
    rows = [
        ("Parsers — PCAP (NAS/NGAP/RRC/F1AP/E1AP/XnAP), Stats (srsRAN/OAI/NIST/Parquet), KPI (Excel/CSV + catalogue)", "Done"),
        ("Detection — 18-detector ensemble (6 PCAP + 6 KPI + 6 Stats)", "Done"),
        ("Event Router — cross-source correlation (cell, category, time window)", "Done"),
        ("Prediction — Prophet + LSTM, 4-hour-ahead forecast", "Done"),
        ("RAG + LLM Explanation — FAISS+MiniLM -> Ollama phi3:mini, rule-based fallback", "Done"),
        ("MLOps — feedback store + nightly retrainer", "Done"),
        ("Interfaces — Streamlit dashboard, FastAPI REST API, CLI orchestrator", "Done"),
        ("Validation — 56 automated tests, all passing", "Done"),
        ("Journal paper draft (Phase I)", "Pending — Review 3"),
    ]
    col_x = [0.3, 7.6]
    col_w = [7.0, 2.1]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.75, col_w[i], 0.3, h, 12, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.12 + r * 0.48
        color = GREEN if row[1] == "Done" else CORAL
        textbox(slide, col_x[0], top, col_w[0], 0.46, row[0], 10, WHITE, wrap=True)
        textbox(slide, col_x[1], top, col_w[1], 0.46, row[1], 10, color, bold=True)

    textbox(slide, 0.3, 5.1, 9.2, 0.4,
            "8 of 9 items done -> Phase 1 code implementation exceeds the "
            "80% target for Review 2; journal paper write-up is the "
            "remaining item for Review 3 (Final Phase).",
            10, GREY, italic=True, wrap=True)


# ── Slide 10: References ────────────────────────────────────────────────

def s_references(prs):
    slide = new_slide(prs)
    slide_title(slide, "References")
    refs = [
        "3GPP TS 24.501 — Non-Access-Stratum (NAS) protocol for 5G System",
        "3GPP TS 38.413 — NG Application Protocol (NGAP)",
        "3GPP TS 38.331 — Radio Resource Control (RRC) protocol specification",
        "3GPP TS 38.473 — F1 Application Protocol (F1AP)",
        "3GPP TS 38.463 — E1 Application Protocol (E1AP)",
        "3GPP TS 38.423 — Xn Application Protocol (XnAP)",
        "3GPP TS 38.913 — Study on scenarios and requirements for next generation access technologies",
        "scikit-learn — Isolation Forest, One-Class SVM, LOF, Elliptic Envelope (Pedregosa et al., 2011)",
        "Hochreiter & Schmidhuber — Long Short-Term Memory (LSTM), Neural Computation 1997",
        "Taylor & Letham — Forecasting at Scale (Prophet), 2018",
        "Johnson, Douze, Jegou — Billion-scale similarity search with GPUs (FAISS), 2017",
        "Reimers & Gurevych — Sentence-BERT (all-MiniLM-L6-v2), EMNLP 2019",
        "Ollama project — local LLM runtime (phi3:mini, llama3.2, mistral)",
    ]
    for i, r in enumerate(refs):
        top = 0.75 + i * 0.35
        textbox(slide, 0.4, top, 9.2, 0.32, f"[{i+1}]  {r}", 11, WHITE, wrap=True)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    s_title(prs)
    s_abstract(prs)
    s_algo_architecture(prs)
    s_algo_pcap(prs)
    s_algo_kpi(prs)
    s_algo_pipeline(prs)
    s_contribution(prs)
    s_results(prs)
    s_code_status(prs)
    s_references(prs)

    out = Path("docs/Phase1_Review2.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
