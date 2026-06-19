"""
Generate the Phase I Second Review PPT — Unified Telecom Analyzer
Academic-review structure: Abstract, Intro, Literature Survey, Objectives,
Proposed System, Architecture, Methodology (parsing/detection/correlation/
prediction/RAG-LLM/MLOps), Implementation, Results, Performance Evaluation,
Model Comparison, Test Results, Conclusion, References.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x2D, 0xC6, 0x53)
ORANGE   = RGBColor(0xFF, 0x94, 0x00)
RED      = RGBColor(0xFF, 0x4D, 0x4D)
GREY     = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE   = RGBColor(0xA0, 0x60, 0xFF)
PINK     = RGBColor(0xFF, 0x4D, 0x9F)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size, color, line_h=0.32):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_title(slide, text):
    textbox(slide, 0.3, 0.12, 9.4, 0.5, text, 22, ACCENT, bold=True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    return slide


# ── Slide 1: Title ────────────────────────────────────────────────────

def s_title(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 0.9, 9.0, 0.9,
            "Unified Telecom Analyzer", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.95, 9.0, 0.5,
            "Multi-Modal Anomaly Detection & Explanation Framework for 5G Networks",
            18, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.5, 9.0, 0.4,
            "PHASE I — SECOND REVIEW  (Research Objective 1)",
            16, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.1, 9.0, 0.4,
            "M.Tech Data Science  |  PES University, Bangalore — Great Learning",
            13, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.55, 9.0, 0.4,
            "Candidate: Vikas  |  Repository: telecom-analyzer",
            12, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.6, 9.0, 0.4,
            "Python · pyshark · scikit-learn · PyTorch · Prophet · FAISS · Ollama · Streamlit · FastAPI",
            11, GREY, align=PP_ALIGN.CENTER)


# ── Slide 2: Abstract ─────────────────────────────────────────────────

def s_abstract(prs):
    slide = new_slide(prs)
    slide_title(slide, "Abstract")
    text = (
        "5G core and radio-access networks generate enormous volumes of signalling "
        "(PCAP), performance counters (DU/CU stats), and KPI time-series, making "
        "manual fault diagnosis slow and error-prone. This project presents a "
        "unified, multi-modal anomaly-detection framework that ingests PCAP traces "
        "(NAS, NGAP, RRC, F1AP, E1AP, XnAP), DU/CU L1/L2 statistics (srsRAN/OAI/"
        "NIST/Parquet), and gNB KPI exports (Excel/CSV) through a single pipeline.\n\n"
        "An 18-detector ensemble (6 methods x 3 data domains) flags anomalies, "
        "a cross-source event router correlates related anomalies across PCAP/KPI/"
        "Stats, a Prophet+LSTM prediction layer forecasts anomalies up to 4 hours "
        "ahead, and a Retrieval-Augmented Generation (RAG) pipeline grounded in "
        "3GPP specifications (TS 24.501, 38.413, 38.331, 38.473, 38.463, 38.423) "
        "produces human-readable root-cause explanations via a locally-hosted LLM "
        "(Ollama phi3:mini), with a deterministic rule-based fallback. An MLOps "
        "feedback loop lets engineers label detections as true/false positives, "
        "feeding a nightly retrainer that tunes detector sensitivity automatically.\n\n"
        "The system is exposed via a Streamlit dashboard, a FastAPI REST API, and "
        "a CLI orchestrator, and is validated by 56 automated tests covering "
        "parsing, detection, correlation, prediction, and feedback."
    )
    textbox(slide, 0.4, 0.75, 9.2, 4.6, text, 12, WHITE, wrap=True)


# ── Slide 3: Introduction / Problem Statement ────────────────────────

def s_intro(prs):
    slide = new_slide(prs)
    slide_title(slide, "Introduction & Problem Statement")
    bullets(slide, 0.4, 0.75, 9.2, 4.6, [
        "5G networks expose multiple heterogeneous data sources: control-plane "
        "signalling (PCAP), physical/MAC-layer counters (DU/CU stats), and "
        "network-level KPIs (Excel/CSV exports).",
        "Operators today analyze each source in isolation, using separate tools "
        "and manual thresholding — this is slow, inconsistent, and misses "
        "anomalies that only appear when sources are correlated.",
        "Single-method anomaly detectors (e.g., only Isolation Forest, or only "
        "fixed thresholds) miss either global statistical outliers, local density "
        "anomalies, or sequence/order violations — no one method covers all cases.",
        "Root-cause explanation is typically manual: an engineer must read 3GPP "
        "specs and cross-reference logs by hand.",
        "Goal: build ONE pipeline that parses all three data types, runs an "
        "ensemble of complementary detectors per domain, correlates findings "
        "across sources, predicts emerging issues ahead of time, and explains "
        "each anomaly with spec-grounded LLM reasoning — with a feedback loop "
        "so detection quality improves over time.",
    ], 12, WHITE, line_h=0.5)


# ── Slide 4: Literature Survey ───────────────────────────────────────

def s_literature(prs):
    slide = new_slide(prs)
    slide_title(slide, "Literature Survey — Summary")
    rows = [
        ("Rule/threshold-based NMS alarms",
         "Vendor EMS/NMS dashboards", "Simple, interpretable, instant",
         "High false positives/negatives; static thresholds don't adapt; no cross-source view"),
        ("Classical ML anomaly detection",
         "Isolation Forest, One-Class SVM, LOF, Elliptic Envelope (literature on network anomaly detection)",
         "Captures statistical outliers without labels",
         "Single method misses other anomaly shapes (local vs global vs sequential)"),
        ("Deep sequence models",
         "LSTM/Autoencoder-based intrusion & fault detection studies",
         "Captures temporal/order dependencies",
         "Needs more data, less interpretable on its own"),
        ("Time-series forecasting",
         "Prophet, ARIMA, LSTM forecasting literature",
         "Early-warning before SLA breach",
         "Rarely integrated with live anomaly detection pipelines"),
        ("LLM + RAG for fault diagnosis",
         "Recent telecom AIOps / RAG-over-specs research",
         "Human-readable, spec-grounded explanations",
         "Cloud LLMs raise data-privacy concerns for operator data"),
    ]
    textbox(slide, 0.4, 0.7, 9.2, 0.3,
            "Gap identified: no open framework combines multi-domain parsing + "
            "multi-method ensemble + cross-source correlation + forecasting + "
            "local LLM explanation + feedback-driven retraining.",
            11, ORANGE, italic=True, wrap=True)

    headers = ["Approach", "Reference / Context", "Strength", "Limitation"]
    top = 1.15
    col_x = [0.3, 2.6, 5.0, 6.9]
    col_w = [2.2, 2.3, 1.8, 2.5]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], top, col_w[i], 0.3, h, 11, ACCENT, bold=True)
    for r, (a, b, c, d) in enumerate(rows):
        ytop = top + 0.35 + r * 0.78
        for i, val in enumerate((a, b, c, d)):
            textbox(slide, col_x[i], ytop, col_w[i], 0.75, val, 9, WHITE, wrap=True)


# ── Slide 5: Research Objective & Scope ──────────────────────────────

def s_objectives(prs):
    slide = new_slide(prs)
    slide_title(slide, "Research Objective 1 — Scope of Phase I")
    bullets(slide, 0.4, 0.75, 9.2, 4.6, [
        "Objective: Design and implement a unified pipeline that parses 5G PCAP, "
        "DU/CU stats, and KPI data, detects anomalies with a multi-method ensemble "
        "per domain, correlates anomalies across sources, forecasts future "
        "anomalies, and explains root causes using a 3GPP-grounded local LLM.",
        "Sub-objective 1 — Parsing: full 3GPP protocol stack decode (NAS, NGAP, "
        "RRC, F1AP, E1AP, XnAP) + DU/CU stats (srsRAN/OAI/NIST/Parquet) + KPI "
        "(Excel/CSV vs 3GPP catalogue).",
        "Sub-objective 2 — Detection: 18-detector ensemble (6 PCAP + 6 KPI + 6 "
        "Stats), each returning severity, rationale, and a comparison matrix.",
        "Sub-objective 3 — Correlation & Prediction: cross-source event router; "
        "Prophet + LSTM forecasting 4h ahead.",
        "Sub-objective 4 — Explanation: FAISS+MiniLM RAG over 3GPP specs -> "
        "Ollama phi3:mini explanation with rule-based fallback.",
        "Sub-objective 5 — MLOps: feedback store + nightly retrainer that adjusts "
        "detector parameters based on observed false-positive rate.",
        "Sub-objective 6 — Interfaces & validation: Streamlit dashboard, FastAPI "
        "REST API, CLI orchestrator, 56 automated tests.",
    ], 11.5, WHITE, line_h=0.45)


# ── Slide 6: Proposed System / Architecture (5 layers) ──────────────

def s_architecture(prs):
    slide = new_slide(prs)
    slide_title(slide, "Proposed System — 5-Layer Architecture")
    layers = [
        ("INPUT",
         "PCAP (.pcap/.pcapng)  ·  DU/CU Stats (srsRAN/OAI/NIST CSV/Parquet)  ·  KPI Time-series (Excel/CSV)",
         ACCENT, 0.75),
        ("PARSING — TS 24.501 / 38.413 / 38.331 / 38.473 / 38.463 / 38.423",
         "NAS · NGAP · RRC · F1AP · E1AP · XnAP · KPI Parser · Stats Parser (srsRAN/OAI/NIST)",
         GREEN, 1.55),
        ("DETECTION — 18 detectors total",
         "PCAP: IF · Statistical · OC-SVM · LOF · Elliptic Env. · LSTM-AE   "
         "KPI: Threshold · Peer Comp. · Trend · IQR · CUSUM · Bollinger   "
         "Stats: same 6 methods on L1/L2 counters",
         ORANGE, 2.45),
        ("CORRELATION + PREDICTION + RAG/LLM + FEEDBACK",
         "Event Router (cross-source correlation)  |  Prophet + LSTM 4h-ahead forecast  |  "
         "FAISS+MiniLM RAG -> Ollama phi3:mini explanation  |  Feedback store -> nightly retrainer",
         PURPLE, 3.35),
        ("OUTPUT",
         "Streamlit Dashboard (http://localhost:8501)  ·  FastAPI REST API (:8000)  ·  CLI Pipeline Runner",
         PINK, 4.2),
    ]
    for label, items, color, top in layers:
        textbox(slide, 0.3, top, 2.7, 0.3, label, 9, color, bold=True)
        textbox(slide, 0.3, top + 0.28, 9.2, 0.5, items, 10, WHITE, wrap=True)


# ── Slide 7: Block / Data-flow diagram ───────────────────────────────

def s_dataflow(prs):
    slide = new_slide(prs)
    slide_title(slide, "End-to-End Data Flow")
    flow = [
        "Streamlit Dashboard / FastAPI REST API / CLI",
        "Orchestrator  (src/orchestrator/pipeline.py) — run_pcap_pipeline · run_kpi_pipeline · run_stats_pipeline",
        "Parsers — PCAP (NAS/NGAP/RRC/F1AP/E1AP/XnAP) | Stats (srsRAN/OAI/NIST/Parquet) | KPI (Excel/CSV + catalogue)",
        "Detection Ensemble — 18 detectors (IF, OC-SVM, LOF, Elliptic Env., Statistical, LSTM-AE x3 domains)",
        "Event Router (cross-source correlation)   +   Prediction Layer (Prophet + LSTM, 4h ahead)",
        "RAG (FAISS over 38 3GPP-spec chunks) -> Ollama phi3:mini LLM explanation  (rule-based fallback)",
        "Feedback Store (JSONL) -> Nightly Retrainer -> updates detector thresholds/contamination",
    ]
    colors = [ACCENT, GREEN, GREEN, ORANGE, PURPLE, PURPLE, PINK]
    for i, (line, c) in enumerate(zip(flow, colors)):
        top = 0.8 + i * 0.62
        textbox(slide, 0.4, top, 0.4, 0.5, f"{i+1}", 16, c, bold=True)
        textbox(slide, 0.95, top, 8.7, 0.55, line, 11, WHITE, wrap=True)
        if i < len(flow) - 1:
            textbox(slide, 0.5, top + 0.42, 0.4, 0.2, "↓", 12, GREY)


# ── Slide 8: Protocols Parsed ────────────────────────────────────────

def s_protocols(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — 3GPP Protocols Parsed (Full 5G-NR Stack)")
    protocols = [
        ("NAS — TS 24.501", "Registration · Auth · Security Mode · PDU Session Establishment · Deregistration", "UE <-> AMF (N1)"),
        ("NGAP — TS 38.413", "InitialContextSetup · PDUSession Setup/Release · UEContextRelease · Handover · Paging", "gNB-CU <-> AMF (N2)"),
        ("RRC — TS 38.331", "Setup · Reestablishment · Reconfiguration · Release · UE Capability · Measurement", "UE <-> gNB (Uu)"),
        ("F1AP — TS 38.473", "F1 Setup · UE Context Setup/Release · DL/UL RRC Transfer · Initial UL RRC", "gNB-DU <-> gNB-CU-CP (F1)"),
        ("E1AP — TS 38.463", "CU-UP Setup · Bearer Context Setup/Modify/Release · Data Notification", "gNB-CU-CP <-> gNB-CU-UP (E1)"),
        ("XnAP — TS 38.423", "Xn Setup · Handover Request · UE Context Release · SN Addition (MR-DC)", "gNB <-> gNB (Xn)"),
    ]
    textbox(slide, 0.3, 0.65, 9.2, 0.3,
            "pyshark used for real captures; raw-discriminator fallback decodes synthetic test PCAPs without a live dissector.",
            10, GREY, wrap=True)
    for i, (proto, procs, iface) in enumerate(protocols):
        top = 1.0 + i * 0.66
        textbox(slide, 0.3, top, 2.4, 0.3, proto, 12, GREEN, bold=True)
        textbox(slide, 2.8, top, 5.1, 0.45, procs, 10, WHITE, wrap=True)
        textbox(slide, 8.0, top, 1.8, 0.3, iface, 9, GREY, align=PP_ALIGN.RIGHT)


# ── Slide 9: PCAP Detection Ensemble ─────────────────────────────────

def s_pcap_detection(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — PCAP Anomaly Detection (6-Detector Ensemble)")
    detectors = [
        ("1", "Isolation Forest", "Unsupervised / tree", "Anomalies isolate in fewer tree splits. No distribution assumption. O(n log n). Best for global outliers."),
        ("2", "Statistical (Threshold+Cascade)", "Rule-based / domain", "Encodes 3GPP SLA directly. 100% interpretable. Catches congestion chains, timeout storms, NAS->NGAP cascades."),
        ("3", "One-Class SVM", "Kernel / boundary", "Non-linear normal region in RKHS. Catches anomalies near but outside the normal boundary."),
        ("4", "LOF", "Density / local", "Local Outlier Factor — catches cells 3-sigma below their local peer cluster even when globally within range."),
        ("5", "Elliptic Envelope", "Mahalanobis / Gaussian", "Fastest interpretable baseline. Fits multivariate Gaussian; sanity-checks all ML methods."),
        ("6", "LSTM Autoencoder", "Deep learning / sequence", "ONLY method that catches procedure-ORDER violations (e.g., PDU Session before Authentication)."),
    ]
    for i, (num, name, dtype, why) in enumerate(detectors):
        top = 0.78 + i * 0.66
        textbox(slide, 0.3, top, 0.25, 0.42, num, 14, ACCENT, bold=True)
        textbox(slide, 0.6, top, 2.3, 0.42, name, 11, WHITE, bold=True)
        textbox(slide, 2.95, top, 1.7, 0.42, dtype, 10, ORANGE)
        textbox(slide, 4.7, top, 5.0, 0.55, why, 10, GREY, wrap=True)


# ── Slide 10: KPI/Stats Detection Ensemble ───────────────────────────

def s_kpi_detection(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — KPI & L1/L2 Stats Detection (6-Method Ensemble)")
    methods = [
        ("1", "Threshold Violation", "Rule-based / point", "Per-row check vs warning/critical thresholds. Encodes 3GPP/ITU-T SLA directly."),
        ("2", "Peer Comparison", "Statistical / cross-cell", "Z-score vs fleet mean. Catches underperforming cells still within threshold but 3-sigma below peers."),
        ("3", "Trend (Linear Regression)", "Temporal / slope", "Slope > 0.5 unit/hr = actionable degradation even if current value looks fine. Early-warning signal."),
        ("4", "IQR (Tukey Fence)", "Robust / distribution-free", "No Gaussian assumption. Works on skewed KPIs — PRB right-skewed, CQI left-skewed during interference."),
        ("5", "CUSUM", "Sequential / change-point", "Accumulates small deviations. Catches RRC SR 99.5%->98.2% drift 5+ hours before threshold fires."),
        ("6", "Bollinger Bands", "Rolling envelope / burst", "Catches single-timestamp spikes invisible to trend methods. Rolling window gives local context."),
    ]
    textbox(slide, 0.3, 0.65, 9.2, 0.3,
            "Same 6 methods applied to KPI time-series (Excel/CSV) AND DU/CU L1/L2 statistics (srsRAN · OAI · NIST) — 12 of the 18 detectors.",
            10, GREY, wrap=True)
    for i, (num, name, dtype, why) in enumerate(methods):
        top = 1.0 + i * 0.66
        textbox(slide, 0.3, top, 0.25, 0.42, num, 14, ACCENT, bold=True)
        textbox(slide, 0.6, top, 2.3, 0.42, name, 11, WHITE, bold=True)
        textbox(slide, 2.95, top, 1.9, 0.42, dtype, 10, ORANGE)
        textbox(slide, 4.9, top, 4.8, 0.55, why, 10, GREY, wrap=True)


# ── Slide 11: Event Correlation ──────────────────────────────────────

def s_correlation(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — Cross-Source Event Correlation")
    bullets(slide, 0.4, 0.8, 9.2, 4.4, [
        "src/orchestrator/event_router.py tags every detector output with its "
        "source (pcap / kpi / stats) and ingests it into a unified event log.",
        "Correlation rule: anomalies on the same cell/gNB within a configurable "
        "time window are linked — e.g. a PCAP handover-failure event + a KPI "
        "handover-success-rate dip on the same cell/time window.",
        "Output: a 'correlated_events' view in the dashboard showing which "
        "anomalies across domains likely share a root cause, plus a top-cells "
        "summary ranking cells by total correlated anomaly count.",
        "Why it matters: a single-domain tool would report two unrelated alerts; "
        "this layer turns them into one actionable incident.",
    ], 12, WHITE, line_h=0.55)


# ── Slide 12: Prediction Layer ───────────────────────────────────────

def s_prediction(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — Prediction Layer (4-Hour-Ahead Forecast)")
    textbox(slide, 0.3, 0.7, 4.6, 0.3, "Prophet — Bayesian structural time-series", 13, PURPLE, bold=True)
    bullets(slide, 0.5, 1.05, 4.4, 3.0, [
        "Handles seasonality, trends, missing data",
        "Best for KPIs with daily/weekly cycles (PRB, throughput)",
        "Produces confidence intervals for forecast",
    ], 11, WHITE)

    textbox(slide, 5.0, 0.7, 4.6, 0.3, "LSTM Sequence Model — nonlinear shifts", 13, ORANGE, bold=True)
    bullets(slide, 5.2, 1.05, 4.4, 3.0, [
        "Catches abrupt non-seasonal shifts (HARQ NACK, BLER)",
        "Learns temporal dependencies beyond fixed seasonality",
        "Complements Prophet for irregular L1/L2 counters",
    ], 11, WHITE)

    textbox(slide, 0.3, 3.3, 9.2, 0.3, "Output contract", 13, GREEN, bold=True)
    bullets(slide, 0.5, 3.65, 9.0, 1.6, [
        "Predicted anomalies tagged state='predicted' with a lead_time_h field, "
        "using the SAME schema as live detections — dashboard and API treat "
        "predicted and live anomalies uniformly.",
        "Dashboard exposes a Prediction Layer panel with a horizon slider (1-4h).",
    ], 11, WHITE)


# ── Slide 13: RAG + LLM ───────────────────────────────────────────────

def s_rag_llm(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — RAG + LLM Root-Cause Explanation")
    flow = [
        ("1  Anomaly input", "procedure + severity + failure_causes + evidence dict", "from PCAP / KPI / Stats detector"),
        ("2  Query build", "natural-language query from anomaly fields", "_build_query()"),
        ("3  Retrieve", "FAISS cosine search -> top-5 3GPP spec chunks", "all-MiniLM-L6-v2 (80MB), 38 curated chunks"),
        ("4  Prompt", "anomaly + retrieved spec context -> structured JSON prompt", "_build_prompt() | temperature=0.2"),
        ("5  LLM call", "Ollama: phi3:mini -> llama3.2:1b -> mistral:7b (preference chain)", "Runs fully LOCAL — no cloud API, no data leaves premises"),
        ("6  Rule-based fallback", "template explanation built from retrieved chunks", "fires automatically if Ollama unavailable"),
        ("7  Output", "hypothesis · severity · citations (spec section) · investigation_hints", "dashboard + REST API response"),
    ]
    for i, (step, desc, note) in enumerate(flow):
        top = 0.8 + i * 0.58
        textbox(slide, 0.3, top, 1.4, 0.42, step, 12, ACCENT, bold=True)
        textbox(slide, 1.8, top, 4.8, 0.42, desc, 11, WHITE)
        textbox(slide, 6.7, top, 3.0, 0.42, note, 9, GREY, wrap=True)
    textbox(slide, 0.3, 5.0, 9.4, 0.28,
            "Knowledge base: 38 curated chunks — TS 24.501 · 38.413 · 38.331 · 38.473 · 38.463 · 38.423 · 38.913",
            10, ORANGE)


# ── Slide 14: MLOps Feedback Loop ────────────────────────────────────

def s_mlops(prs):
    slide = new_slide(prs)
    slide_title(slide, "Methodology — MLOps Feedback Loop")
    textbox(slide, 0.3, 0.72, 4.5, 0.3, "Feedback Loop", 14, GREEN, bold=True)
    bullets(slide, 0.5, 1.05, 4.3, 3.0, [
        "Engineer marks each anomaly: Correct / False Positive / Uncertain",
        "Feedback stored to JSONL log (src/feedback/store.py)",
        "Nightly retrainer reads FP rate per detector",
        "Adjusts contamination (IF/SVM/LOF/EE), sensitivity (Statistical), "
        "threshold (LSTM)",
        "Formula: new_param = base + alpha x (fp_rate - target_fp_rate), "
        "clamped to [0.02, 0.40]",
        "Scheduled via cron: make retrain (scripts/nightly_retrain.py)",
    ], 11, WHITE, line_h=0.45)

    textbox(slide, 5.0, 0.72, 4.7, 0.3, "Why it matters", 14, PURPLE, bold=True)
    bullets(slide, 5.2, 1.05, 4.5, 3.0, [
        "Closes the loop between detection and ground truth without retraining "
        "from scratch every time",
        "Detector sensitivity self-corrects toward the operator's observed "
        "false-positive tolerance",
        "Tuned parameters persisted to data/models/retrain_config.json and "
        "picked up automatically on the next pipeline run",
    ], 11, WHITE, line_h=0.5)


# ── Slide 15: Interfaces / Implementation ────────────────────────────

def s_interfaces(prs):
    slide = new_slide(prs)
    slide_title(slide, "Implementation — Interfaces")
    interfaces = [
        ("Streamlit Dashboard", "make dev -> http://localhost:8501",
         ["Upload PCAP, Stats CSV, or KPI Excel in one interface",
          "6 protocol tabs, failure-cause breakdown, IE inspector",
          "18-detector method-comparison matrix with rationale expanders",
          "LLM explanation with 3GPP spec citations",
          "Prediction panel, feedback buttons, event-router correlation view"]),
        ("FastAPI REST API", "make api -> http://localhost:8000/docs",
         ["POST /analyze/pcap — upload PCAP, returns full anomaly JSON",
          "POST /analyze/kpi — upload KPI file, returns anomaly JSON",
          "GET /analyze/{job_id} — retrieve cached result",
          "GET /health — Ollama status, model, version"]),
        ("CLI Pipeline Runner",
         "python -m src.orchestrator.pipeline --input file.pcap [--no-llm] [--output report.json]",
         ["Auto-detects PCAP / KPI / Stats from file extension",
          "--no-llm skips LLM for fast/CI runs",
          "--output writes full JSON report for batch processing"]),
    ]
    top = 0.72
    for name, cmd, items in interfaces:
        textbox(slide, 0.3, top, 9.2, 0.3, f"{name}  —  {cmd}", 12, GREEN, bold=True)
        for j, b in enumerate(items):
            textbox(slide, 0.7, top + 0.32 + j * 0.28, 9.0, 0.28, f"• {b}", 10, WHITE)
        top += 0.32 + len(items) * 0.28 + 0.18


# ── Slide 16: Tech Stack ──────────────────────────────────────────────

def s_techstack(prs):
    slide = new_slide(prs)
    slide_title(slide, "Implementation — Technology Stack")
    rows = [
        ("Language / Core", "Python 3.13, pandas, numpy"),
        ("PCAP Parsing", "pyshark (tshark), raw-discriminator fallback"),
        ("ML — PCAP detectors", "scikit-learn (Isolation Forest, One-Class SVM, LOF, Elliptic Envelope), PyTorch (LSTM-AE)"),
        ("ML — KPI/Stats detectors", "numpy/pandas/scipy (Threshold, Peer Z-score, Linear Regression Trend, IQR, CUSUM, Bollinger)"),
        ("Forecasting", "Prophet (seasonal/trend), PyTorch LSTM (nonlinear)"),
        ("RAG", "FAISS (vector index), sentence-transformers all-MiniLM-L6-v2"),
        ("LLM", "Ollama — phi3:mini (primary), llama3.2:1b / mistral:7b (fallback chain)"),
        ("Feedback / MLOps", "JSONL feedback log, nightly retrainer cron job"),
        ("Interfaces", "Streamlit (dashboard), FastAPI (REST API), argparse CLI"),
        ("Testing", "pytest — 56 tests (unit + integration)"),
    ]
    for i, (k, v) in enumerate(rows):
        top = 0.75 + i * 0.43
        textbox(slide, 0.3, top, 2.6, 0.4, k, 11, ACCENT, bold=True)
        textbox(slide, 3.0, top, 6.5, 0.4, v, 10, WHITE, wrap=True)


# ── Slide 17: Experimental Results ───────────────────────────────────

def s_results(prs):
    slide = new_slide(prs)
    slide_title(slide, "Experimental Results")
    bullets(slide, 0.4, 0.78, 9.2, 1.7, [
        "Test data generated synthetically to cover known-good and injected "
        "failure scenarios per data type (PCAP, KPI, Stats), so detector "
        "accuracy can be verified against ground truth.",
        "All 18 detectors return severity (Critical/High/Medium/Low) and a "
        "human-readable rationale; results below are from actual end-to-end "
        "pipeline runs (--no-llm) on the current datasets.",
    ], 11.5, WHITE, line_h=0.4)

    headers = ["Dataset", "Size", "Anomalies (ensemble)", "Key findings"]
    rows = [
        ("ue_attach_10.pcap", "199 pkts / 14 procedures", "12",
         "NAS_Registration & XnAP_HandoverRequest flagged by 4-6 of 6 detectors (IF, OC-SVM, LOF, LSTM-AE)"),
        ("ue_handover_20.pcap", "108 pkts / 6 procedures", "6",
         "NGAP_HandoverPreparation, XnAP_HandoverRequest, RRC_Reconfiguration flagged by IF/OC-SVM/LSTM-AE"),
        ("kpi_10hr_sample.csv", "6000 rows x 25 KPIs", "1486",
         "Threshold 1057, Bollinger 139, CUSUM 176, IQR 97, Peer 17 — congestion/outage/drift injected & recovered"),
        ("srsran_stats.csv", "960 rows x 15 cols", "29",
         "Trend detector (29) — slow drift captured in L1/L2 counters; other 5 methods clean"),
    ]
    col_x = [0.3, 2.3, 4.2, 5.9]
    col_w = [1.9, 1.8, 1.6, 3.3]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 2.55, col_w[i], 0.3, h, 11, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 2.92 + r * 0.62
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.6, val, 9.5, WHITE, wrap=True)

    textbox(slide, 0.3, 5.55, 9.2, 0.3,
            "Predicted anomalies (Prophet/LSTM) were verified to use the same output "
            "schema as live detections (state='predicted', lead_time_h).",
            10, GREY, wrap=True)


# ── Slide 18: Performance Evaluation ─────────────────────────────────

def s_performance(prs):
    slide = new_slide(prs)
    slide_title(slide, "Performance Evaluation")
    headers = ["Pipeline", "Input size", "Measured latency (--no-llm)", "Notes"]
    rows = [
        ("PCAP (attach)", "199 pkts / 14 procedures", "6.44 s", "Includes LSTM-AE training (~3.6s of total); IF/OC-SVM/LOF/EE near-instant"),
        ("PCAP (handover)", "108 pkts / 6 procedures", "1.40 s", "Fewer procedures -> Elliptic Envelope skipped (min-sample guard)"),
        ("KPI", "6000 rows x 25 cols", "11.61 s", "CUSUM (~0.5s) + Bollinger (~6.6s) dominate; Threshold/Peer/Trend near-instant"),
        ("Stats (L1/L2)", "960 rows x 15 cols", "1.05 s", "Same 6-method ensemble, smaller feature set"),
        ("End-to-end (with LLM)", "single anomaly", "+1-3 s", "Ollama phi3:mini local inference; rule-based fallback is instant"),
    ]
    col_x = [0.3, 2.6, 4.7, 6.9]
    col_w = [2.2, 1.9, 2.0, 2.6]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.78, col_w[i], 0.3, h, 12, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.15 + r * 0.78
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.74, val, 10, WHITE, wrap=True)

    textbox(slide, 0.3, 5.05, 9.2, 0.6,
            "All measurements on a single 16GB-RAM WSL2 dev machine, CPU-only "
            "(no GPU), from actual pipeline runs on data/raw/* sample files. "
            "LLM step has an instant deterministic fallback.",
            10, GREY, wrap=True)


# ── Slide 19: Model / Method Comparison Matrix ───────────────────────

def s_model_comparison(prs):
    slide = new_slide(prs)
    slide_title(slide, "Model / Method Comparison Matrix")
    headers = ["Detector family", "Catches", "Misses", "Interpretability", "Speed"]
    rows = [
        ("Isolation Forest", "Global statistical outliers", "Local/density anomalies", "Medium", "Fast"),
        ("One-Class SVM", "Non-linear boundary anomalies", "Sequence/order issues", "Low", "Medium"),
        ("LOF", "Local density anomalies (peer-relative)", "Global trend shifts", "Medium", "Medium"),
        ("Elliptic Envelope", "Gaussian-distribution outliers", "Non-Gaussian / skewed data", "High", "Fastest"),
        ("Statistical / Threshold", "Known SLA breaches, cascades", "Novel/unseen failure patterns", "Highest", "Fastest"),
        ("LSTM Autoencoder", "Procedure-order / sequence violations", "Single-point statistical outliers", "Low", "Slowest"),
        ("Trend / CUSUM", "Slow drift, change-points", "Sudden one-off spikes", "Medium", "Fast"),
        ("Bollinger Bands", "Burst spikes vs rolling baseline", "Slow long-term drift", "Medium", "Fast"),
        ("IQR (Tukey)", "Skewed-distribution outliers", "Gradual trend shifts", "High", "Fastest"),
    ]
    col_x = [0.3, 2.5, 5.0, 7.5, 8.7]
    col_w = [2.1, 2.4, 2.4, 1.1, 1.0]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.72, col_w[i], 0.3, h, 11, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.08 + r * 0.46
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.44, val, 9, WHITE, wrap=True)
    textbox(slide, 0.3, 5.15, 9.2, 0.3,
            "Conclusion: methods are complementary, not redundant — ensemble agreement raises detection confidence and lowers false positives.",
            10, GREEN, italic=True, wrap=True)


# ── Slide 20: Test Results ───────────────────────────────────────────

def s_tests(prs):
    slide = new_slide(prs)
    slide_title(slide, "Validation — 56 Automated Tests, All Passing")
    suites = [
        ("test_walking_skeleton.py (4 tests)", ["parse_pcap_stub keys", "detect_anomalies_stub list", "explain_anomaly_stub citations", "stub pipeline end-to-end"]),
        ("test_real_pipeline.py (16 tests)", ["Real PCAP parser — keys/procedures/schema", "6 PCAP detectors — dict schema", "KPI parser + 6 KPI detectors", "RAG retriever — chunks + relevance", "LLM explainer — rule-based schema", "run_pcap/kpi/pipeline() end-to-end"]),
        ("test_event_router.py (8 tests)", ["ingest + summary + correlated events", "cross-source correlation, top cells, event log"]),
        ("test_feedback_prediction.py (12 tests)", ["feedback save/load/stats", "LSTM + Prophet prediction schema"]),
        ("test_retrainer.py (8 tests)", ["FP rate calculation", "param adjustment formula", "dry-run mode"]),
        ("test_walking_skeleton.py — stats (8 tests)", ["stats parser — srsRAN/OAI/NIST", "6 stats detectors schema"]),
    ]
    columns = [0.3, 5.0]
    for idx, (title, tests) in enumerate(suites):
        col = idx // 3
        row = idx % 3
        x = columns[col]
        top = 0.78 + row * 1.55
        textbox(slide, x, top, 4.5, 0.3, title, 11, ACCENT, bold=True)
        for j, t in enumerate(tests):
            textbox(slide, x + 0.2, top + 0.32 + j * 0.26, 4.3, 0.26, f"OK  {t}", 9.5, GREEN)
    textbox(slide, 0.3, 5.2, 9.2, 0.3, "Run: make test (smoke) / make test-all (full 56-test suite)", 10, GREY)


# ── Slide 21: Conclusion & Contributions ─────────────────────────────

def s_conclusion(prs):
    slide = new_slide(prs)
    slide_title(slide, "Conclusion & Contributions (Phase I — Second Review)")
    bullets(slide, 0.4, 0.78, 9.2, 4.4, [
        "Delivered a single, unified pipeline covering 3 heterogeneous 5G data "
        "domains (PCAP, DU/CU stats, KPI) end to end.",
        "Implemented and validated an 18-detector ensemble (6 methods x 3 "
        "domains), each with documented strengths/limitations and a "
        "cross-method comparison matrix.",
        "Built a cross-source event correlator and a 4-hour-ahead Prophet+LSTM "
        "prediction layer sharing a common anomaly schema.",
        "Built a fully local RAG+LLM explanation pipeline grounded in 6 3GPP "
        "specifications, with deterministic fallback for offline/no-GPU use.",
        "Closed the loop with an engineer feedback store and nightly retrainer "
        "that adapts detector sensitivity from observed false-positive rate.",
        "Exposed all functionality via dashboard, REST API, and CLI; validated "
        "with 56 automated tests, all passing.",
        "Status: Phase I code implementation is complete (exceeds the 80% "
        "target for this review). Third (Final) Review will add the journal "
        "paper write-up and Phase I-vs-Phase II final demo.",
    ], 12, WHITE, line_h=0.5)


# ── Slide 22: References ─────────────────────────────────────────────

def s_references(prs):
    slide = new_slide(prs)
    slide_title(slide, "References")
    refs = [
        "3GPP TS 24.501 — Non-Access-Stratum (NAS) protocol for 5G System",
        "3GPP TS 38.413 — NG Application Protocol (NGAP)",
        "3GPP TS 38.331 — Radio Resource Control (RRC) protocol specification",
        "3GPP TS 38.473 — F1 Application Protocol (F1AP)",
        "3GPP TS 38.463 — E1 Application Protocol (E1AP)",
        "3GPP TS 38.423 — Xn Application Protocol (XnAP)",
        "3GPP TS 38.913 — Study on scenarios and requirements for next generation access technologies",
        "scikit-learn — Isolation Forest, One-Class SVM, LOF, Elliptic Envelope (Pedregosa et al., 2011)",
        "Hochreiter & Schmidhuber — Long Short-Term Memory (LSTM), Neural Computation 1997",
        "Taylor & Letham — Forecasting at Scale (Prophet), 2018",
        "Johnson, Douze, Jegou — Billion-scale similarity search with GPUs (FAISS), 2017",
        "Reimers & Gurevych — Sentence-BERT (all-MiniLM-L6-v2), EMNLP 2019",
        "Ollama project — local LLM runtime (phi3:mini, llama3.2, mistral)",
    ]
    for i, r in enumerate(refs):
        top = 0.75 + i * 0.35
        textbox(slide, 0.4, top, 9.2, 0.32, f"[{i+1}]  {r}", 11, WHITE, wrap=True)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    s_title(prs)
    s_abstract(prs)
    s_intro(prs)
    s_literature(prs)
    s_objectives(prs)
    s_architecture(prs)
    s_dataflow(prs)
    s_protocols(prs)
    s_pcap_detection(prs)
    s_kpi_detection(prs)
    s_correlation(prs)
    s_prediction(prs)
    s_rag_llm(prs)
    s_mlops(prs)
    s_interfaces(prs)
    s_techstack(prs)
    s_results(prs)
    s_performance(prs)
    s_model_comparison(prs)
    s_tests(prs)
    s_conclusion(prs)
    s_references(prs)

    out = Path("docs/Phase1_Second_Review.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
