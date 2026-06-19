"""
Generate "Algorithm & Method Justification" deck — Unified Telecom Analyzer.

Purpose: viva/review reference — for every algorithm used in the pipeline,
explain WHAT it does, WHERE it's used, KEY PARAMS, and WHY it was chosen
over common alternatives. Ends with a Q&A cheat-sheet slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x2D, 0xC6, 0x53)
ORANGE   = RGBColor(0xFF, 0x94, 0x00)
RED      = RGBColor(0xFF, 0x4D, 0x4D)
GREY     = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE   = RGBColor(0xA0, 0x60, 0xFF)
PINK     = RGBColor(0xFF, 0x4D, 0x9F)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size, color, line_h=0.32):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_title(slide, text, sub=None):
    textbox(slide, 0.3, 0.10, 9.4, 0.45, text, 21, ACCENT, bold=True)
    if sub:
        textbox(slide, 0.3, 0.52, 9.4, 0.3, sub, 11, GREY, italic=True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    return slide


def table(slide, headers, rows, col_x, col_w, header_y, row_h, start_y,
          font_size=9.5, header_size=11):
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], header_y, col_w[i], 0.3, h, header_size, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = start_y + r * row_h
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], row_h - 0.04, val, font_size, WHITE, wrap=True)


# ── Slide 1: Title ────────────────────────────────────────────────────

def s_title(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 0.85, 9.0, 0.9,
            "Algorithm & Method Justification", 34, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.55, 9.0, 0.5,
            "Unified Telecom Analyzer — Why this method, and not the alternatives",
            16, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.15, 9.0, 0.4,
            "Viva / Review Reference — Phase I",
            14, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.7, 9.0, 0.4,
            "M.Tech Data Science  |  PES University, Bangalore — Great Learning",
            12, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.1, 9.0, 0.4,
            "Candidate: Vikas  |  Repository: telecom-analyzer",
            11, GREY, align=PP_ALIGN.CENTER)
    bullets(slide, 1.0, 3.9, 8.0, 1.4, [
        "For each pipeline stage: what the method is, what data it sees, key hyperparameters",
        "Why it was chosen over standard alternatives (interpretability, data size, CPU-only, offline constraints)",
        "Closing slide = quick-fire Q&A cheat-sheet for the panel",
    ], 12, WHITE, line_h=0.4)


# ── Slide 2: Pipeline overview map ──────────────────────────────────────

def s_overview(prs):
    slide = new_slide(prs)
    slide_title(slide, "Where Each Method Sits in the Pipeline")
    headers = ["Stage", "Method family used", "Core idea"]
    rows = [
        ("1. PCAP parsing", "pyshark (real) + raw byte-discriminator fallback (synthetic)", "Full 3GPP dissection vs lightweight test-traffic decoder"),
        ("2. PCAP anomaly detection", "6-detector ensemble: IF, OC-SVM, LOF, Elliptic Envelope, Rule-based, LSTM-AE", "Tabular + sequential, unsupervised, complementary"),
        ("3. KPI anomaly detection", "6-method ensemble: Threshold, Peer, Trend, IQR, CUSUM, Bollinger", "Absolute / relative / time-horizon complementary checks"),
        ("4. Stats (L1/L2) detection", "Same 6-method ensemble, noise-tuned thresholds", "Reuse with looser fences for PHY/MAC counters"),
        ("5. Cross-source correlation", "Rule-based bucketing by cell_id / procedure-category", "Deterministic, auditable NOC-style cross-check"),
        ("6. Prediction (4h ahead)", "Prophet + lightweight LSTM forecaster", "Decomposition model + nonlinear model, cross-validated"),
        ("7. Root-cause explanation", "FAISS + MiniLM RAG -> Ollama phi3:mini LLM (rule-based fallback)", "Local, offline, small-model + retrieval"),
        ("8. MLOps feedback loop", "Proportional controller on detector thresholds", "Tunes sensitivity from false-positive rate, no retraining"),
    ]
    col_x = [0.3, 2.5, 6.1]
    col_w = [2.1, 3.5, 3.7]
    table(slide, headers, rows, col_x, col_w, header_y=0.78, row_h=0.575, start_y=1.12, font_size=9)


# ── Slide 3: PCAP Anomaly Ensemble ──────────────────────────────────────

def s_pcap_ensemble(prs):
    slide = new_slide(prs)
    slide_title(slide, "PCAP Anomaly Detection — Why a 6-Detector Ensemble?",
                 "Each detector covers a blind spot of the others — chosen for complementary coverage, not redundancy")
    headers = ["Method", "Key params", "Why this, not the alternative"]
    rows = [
        ("Isolation Forest\n(global outliers)", "contamination=0.15\nn_estimators=200", "No distribution assumption, O(n log n) — vs DBSCAN (needs density tuning) or Gaussian methods (assumes normality). Works with very few procedures and zero labels."),
        ("One-Class SVM\n(boundary anomalies)", "nu=0.15, kernel=rbf", "Learns a non-linear boundary IF's axis-aligned splits miss. One-class formulation needs no labeled negative class — vs 2-class SVM which would need labeled anomalies we don't have."),
        ("Local Outlier Factor\n(peer-relative density)", "n_neighbors=5\ncontamination=0.15", "Gives a continuous anomaly score (needed for severity ranking) and needs no epsilon radius — vs DBSCAN's binary cluster/noise label and sensitive eps tuning."),
        ("Elliptic Envelope\n(Mahalanobis distance)", "contamination=0.10\n(tighter)", "Single, statistically-grounded distance metric — most interpretable baseline. Tighter threshold by design, flags only the clearest outliers; skips gracefully on small samples instead of guessing."),
        ("Statistical / Rule-based\n(3GPP SLA rules)", "Fixed thresholds:\nSR<98/95/85%, cascade chains", "Fully interpretable, traceable to 3GPP SLA targets and procedure-order rules (e.g. RRC before NGAP) that NO ML model can encode without labels. Acts as ground-truth sanity check."),
        ("LSTM Autoencoder\n(sequence/order)", "hidden=32, epochs=60\nflag if err > mean+2*std", "Only detector that sees message ORDER, not just feature values — catches a Reconfig before Setup completes. No labels needed; trains in seconds per capture — vs a transformer (needs pretraining) or n-gram models (can't model variable-length deps)."),
    ]
    col_x = [0.3, 2.1, 4.0]
    col_w = [1.75, 1.85, 5.9]
    table(slide, headers, rows, col_x, col_w, header_y=1.0, row_h=0.685, start_y=1.32, font_size=8.8)


# ── Slide 4: KPI Detection Ensemble ─────────────────────────────────────

def s_kpi_ensemble(prs):
    slide = new_slide(prs)
    slide_title(slide, "KPI Anomaly Detection — 6 Complementary Methods",
                 "Spans 3 axes: absolute vs relative, point vs trend, Gaussian vs distribution-free")
    headers = ["Method", "Catches", "Why this, not the alternative"]
    rows = [
        ("Threshold", "Absolute SLA breach (RRC SR<99% etc.)", "Direct encoding of operator/3GPP SLA limits — instant, zero training data. The 'ground truth' the statistical methods are validated against."),
        ("Peer Comparison\n(z-score)", "Relative under-performance vs fleet (>2sigma / >3sigma)", "Catches a cell that's 'fine' on absolute SLA but a major outlier vs peers. Z-score chosen over MAD/robust-scale for simplicity with small cell counts."),
        ("Trend\n(linear regression slope)", "Slow degradation (slope>0.5/hr)", "Proactive vs Threshold's reactive check. Plain OLS slope chosen over STL decomposition — 10-hr sample is too short for seasonal decomposition; slope is directly interpretable."),
        ("IQR (Tukey fence)", "Skewed-distribution outliers (PRB, CQI)", "Non-parametric, no Gaussian assumption — vs z-score which misjudges fences on right/left-skewed telecom KPIs. Config-free."),
        ("CUSUM\n(change-point)", "Sustained gradual drift (99.5% -> 98.2% over 10 intervals)", "Industry-standard NOC method for early degradation warning. Detects collective small deviations individually below any threshold — vs EWMA, simpler and the textbook standard."),
        ("Bollinger Bands", "Short-term spikes/dips (rolling mean +/- 2sigma)", "Fills the gap Trend (long-term) and CUSUM (sustained) leave open — transient PRB spikes / CQI crashes. Borrowed from finance, proven for KPI burst detection."),
    ]
    col_x = [0.3, 2.0, 4.0]
    col_w = [1.65, 1.95, 5.9]
    table(slide, headers, rows, col_x, col_w, header_y=1.0, row_h=0.685, start_y=1.32, font_size=8.8)


# ── Slide 5: Stats / L1L2 Ensemble — what's different ───────────────────

def s_stats_ensemble(prs):
    slide = new_slide(prs)
    slide_title(slide, "L1/L2 Stats Detection — Same Ensemble, Re-tuned for Noise",
                 "Reuses the 6 KPI methods; fences widened and made percentage-based for sub-second PHY/MAC counters")
    headers = ["Method", "KPI version", "Stats (L1/L2) version", "Why the change"]
    rows = [
        ("Severity direction", "Explicit 'direction' field in KPI metadata", "String-matched on metric name (BLER/PRB/HARQ-NACK=high-bad; MCS/SNR/RSRP=low-bad)", "No standard metadata catalogue for raw PHY/MAC counters across srsRAN/OAI/NIST"),
        ("IQR fence", "1.5 x IQR, per-row flag", "3 x IQR, flag only if >5% of points violate", "Millisecond-scale counters are far noisier; a single-sample violation is mostly noise"),
        ("CUSUM", "slack=0.5sigma, h=5sigma", "slack=0.5, thresh=4.0 (more sensitive)", "Tuned for the higher native variance of L1/L2 data"),
        ("Bollinger", "k=2sigma, worst single point", "k=2.5sigma, needs >=3% of points outside band", "Avoid false alarms from PHY-layer jitter"),
        ("Output", "All anomalies from all 6 detectors kept", "De-duplicated by (label, cell_id), keep highest severity", "L1/L2 ensemble would otherwise massively over-report on noisy counters"),
    ]
    col_x = [0.3, 1.7, 3.6, 6.4]
    col_w = [1.4, 1.9, 2.8, 3.1]
    table(slide, headers, rows, col_x, col_w, header_y=1.0, row_h=0.82, start_y=1.32, font_size=9)


# ── Slide 6: Prediction Layer ────────────────────────────────────────────

def s_prediction(prs):
    slide = new_slide(prs)
    slide_title(slide, "Prediction Layer (4-Hour-Ahead Forecast)",
                 "Prophet + LSTM run in parallel and cross-validate each other's forecast")
    headers = ["Model", "Key params", "Why this, not the alternative"]
    rows = [
        ("Prophet\n(additive decomposition)", "interval_width=0.95\nseasonality off (series too short)\nchangepoint_prior=0.1, min 24 rows", "Robust to missing/irregular intervals; gives uncertainty bounds (yhat_upper) used directly for breach-escalation. Chosen over ARIMA — ARIMA needs manual stationarity testing (ADF) and ACF/PACF tuning per cell x KPI, impractical across dozens of combinations. Prophet auto-tunes changepoints and is interpretable for a thesis defense."),
        ("Lightweight LSTM\n(per cell x KPI)", "lookback=12, hidden=32\nepochs=30, Adam lr=1e-2\nautoregressive multi-step", "Captures nonlinear dependencies Prophet's additive model can't (e.g. BLER's nonlinear response to interference). Trains in <1s per series — no pretrained weights/persistence needed, fits an offline student-project MLOps loop."),
        ("Why run BOTH", "Both forecasts compared per cell x KPI", "If Prophet and LSTM agree on a future breach, confidence in the prediction is higher — a simple, explainable ensemble rather than betting on a single forecaster."),
    ]
    col_x = [0.3, 2.3, 4.3]
    col_w = [1.95, 1.95, 5.6]
    table(slide, headers, rows, col_x, col_w, header_y=1.0, row_h=1.35, start_y=1.35, font_size=9.5)


# ── Slide 7: Cross-Source Correlation ────────────────────────────────────

def s_correlation(prs):
    slide = new_slide(prs)
    slide_title(slide, "Cross-Source Event Correlation",
                 "Rule-based bucketing — not a learned model")
    bullets(slide, 0.4, 1.0, 9.2, 2.2, [
        "Method: bucket anomalies by shared key — cell_id where available, else procedure-category prefix "
        "(proc::rrc, proc::nas, proc::ngap) for PCAP events without a cell_id.",
        "Any bucket touched by >= 2 distinct sources (pcap / kpi / stats) gets correlated_sources populated "
        "and correlation_confidence = (number of sources) / 3.",
        "Final ranking: sort by (severity_rank, num_correlated_sources, score) descending — multi-source "
        "agreement pushes an anomaly up the priority list.",
    ], 12, WHITE, line_h=0.55)
    textbox(slide, 0.4, 3.55, 9.2, 0.35, "Why rule-based, not a learned correlation model (clustering / GNN):", 13, ACCENT, bold=True)
    bullets(slide, 0.4, 3.95, 9.2, 1.6, [
        "Only 3 sources and a simple shared key (cell_id) — a learned model would be massive overkill and "
        "impossible to explain to examiners.",
        "Deterministic and auditable: mirrors exactly how a NOC engineer manually cross-checks alarms "
        "('same cell flagged by KPI and Stats = high confidence').",
        "The heuristic confidence score (n_sources/3) is simple to defend and re-tune without any retraining.",
    ], 11.5, WHITE, line_h=0.45)


# ── Slide 8: RAG (Retrieval) ─────────────────────────────────────────────

def s_rag(prs):
    slide = new_slide(prs)
    slide_title(slide, "RAG Retrieval — MiniLM Embeddings + FAISS",
                 "Small, local, exact-search stack sized to a ~36-chunk curated 3GPP knowledge base")
    headers = ["Component", "Choice used", "Why this, not the alternative"]
    rows = [
        ("Embedding model", "all-MiniLM-L6-v2\n(~80 MB, sentence-transformers)", "CPU inference <10ms/query, 80MB download vs 1GB+ for bge-large / OpenAI embeddings. With only ~36 chunks, larger embedding models give negligible quality gain — not worth the offline/CPU cost."),
        ("Vector store", "FAISS IndexFlatIP\n(exact search, normalized -> cosine)", "Exact (not approximate) search is fine for ~36 vectors — no need for IVF/HNSW ANN indexes. FAISS is a local library: no server, no network dependency, fully reproducible — vs Pinecone/Chroma/Weaviate which add infra overhead for no benefit at this scale."),
        ("Knowledge base", "~36 hand-curated chunks across\nTS 24.501, 38.413, 38.331, 38.473, 38.463, 38.423", "3GPP specs are huge and dense; automatic chunking (recursive splitter over raw PDFs) would yield noisy, poorly-bounded excerpts. Manual curation maps each chunk cleanly to one failure-cause/procedure concept -> higher retrieval precision for a small domain-specific corpus."),
        ("Retrieval query", "Concatenation of anomaly type, procedure, layer, evidence, top failure cause; top-k=5", "Query built from structured anomaly fields (not free text) so retrieval is precise even with a small, literal embedding model."),
    ]
    col_x = [0.3, 1.9, 4.3]
    col_w = [1.55, 2.35, 5.55]
    table(slide, headers, rows, col_x, col_w, header_y=1.0, row_h=1.0, start_y=1.35, font_size=9)


# ── Slide 9: LLM Explainer ────────────────────────────────────────────────

def s_llm(prs):
    slide = new_slide(prs)
    slide_title(slide, "Root-Cause Explanation — Local LLM (Ollama / phi3:mini)",
                 "Offline, CPU-only, privacy-preserving — with a deterministic rule-based fallback")
    bullets(slide, 0.4, 1.0, 9.2, 1.7, [
        "Model: Ollama running locally; preference list phi3:mini -> phi3 -> phi3:3.8b -> llama3.2:1b/3b -> "
        "mistral:7b-instruct -> gemma2:2b (first available is auto-selected).",
        "Inference: top-5 RAG chunks injected into prompt, temperature=0.2, num_predict=600 -> structured JSON "
        "(hypothesis, citations, suggested action).",
        "Fallback: if Ollama is unreachable or returns invalid JSON, _rule_based_explanation() builds the same "
        "structured output directly from retrieved chunks + a cause-to-explanation dictionary — system still "
        "works, never crashes.",
    ], 11.5, WHITE, line_h=0.5)
    textbox(slide, 0.4, 3.0, 9.2, 0.35, "Why phi3:mini + Ollama, not GPT-4 / Claude API:", 13, ACCENT, bold=True)
    bullets(slide, 0.4, 3.4, 9.2, 2.0, [
        "Offline requirement: no internet dependency or per-call API cost for repeated batch anomaly runs.",
        "CPU-only feasibility: phi3:mini (3.8B, quantized) runs on a laptop CPU in seconds; GPT-4-class models "
        "need GPU/cloud.",
        "Data privacy: telecom captures can contain operator-sensitive info — local inference avoids sending "
        "data to third-party APIs.",
        "RAG does the heavy lifting — retrieval already supplies the factual 3GPP context, so a small "
        "instruction-tuned model only needs to reformat/synthesize, not 'know' everything.",
        "Graceful degradation (rule-based fallback) is itself a 'production-grade' design point for examiners.",
    ], 11, WHITE, line_h=0.4)


# ── Slide 10: MLOps Retrainer ─────────────────────────────────────────────

def s_retrainer(prs):
    slide = new_slide(prs)
    slide_title(slide, "MLOps Feedback Loop — Proportional Controller, Not Retraining",
                 "Nightly job tunes detector sensitivity from engineer feedback")
    bullets(slide, 0.4, 1.0, 9.2, 2.0, [
        "Trigger: nightly batch reads engineer feedback (data/feedback); needs >=5 rated records, else skipped.",
        "What's adjusted: thresholds/hyperparameters only — contamination/nu for IF/LOF/EE/OC-SVM, sensitivity "
        "for the rule-based detector, threshold_multiplier for LSTM-AE, and z/slope/k/threshold for KPI & "
        "stats detectors. Model WEIGHTS are never touched.",
        "Algorithm: simple P-controller —  new = old +/- ALPHA * (fp_rate - TARGET_FP_RATE), ALPHA=0.30, "
        "TARGET_FP_RATE=0.10, all values clamped to safe ranges (e.g. contamination in [0.02, 0.40]).",
        "Output: data/models/retrain_config.json, read by every detector via load_retrain_config() on the next run.",
    ], 11.5, WHITE, line_h=0.5)
    textbox(slide, 0.4, 3.6, 9.2, 0.35, "Why threshold-tuning, not full model retraining:", 13, ACCENT, bold=True)
    bullets(slide, 0.4, 4.0, 9.2, 1.4, [
        "No large labeled dataset exists to retrain IF / OC-SVM / LOF / LSTM-AE from scratch nightly.",
        "A P-controller update is computationally trivial, immediately explainable ('too many false alarms -> "
        "dial sensitivity down by a fixed step'), and converges smoothly with no catastrophic forgetting.",
        "Mirrors real MLOps practice: unsupervised detectors have no loss tied to 'correctness' — only "
        "downstream alert precision/volume can be measured and fed back.",
    ], 11, WHITE, line_h=0.4)


# ── Slide 11: PCAP Parsing approach ───────────────────────────────────────

def s_parsing(prs):
    slide = new_slide(prs)
    slide_title(slide, "PCAP Parsing — Dual-Path Design")
    headers = ["Path", "Used for", "Approach", "Why"]
    rows = [
        ("pyshark dissection", "Real-world captures", "Full 3GPP protocol stacks (NAS / NGAP / RRC / F1AP / E1AP / XnAP) with deep IE extraction via ie_extractor", "Standards-compliant, handles any real capture without a custom encoder"),
        ("Raw byte-discriminator fallback", "Synthetic test PCAPs from the project's own generators", "Single-byte protocol discriminators (0x7e=NAS, 0x3a=NGAP, 0x2b=RRC, 0x1a=F1AP, ...)", "Avoids full ASN.1 encoding overhead for test traffic — lightweight, fast to generate and parse, still exercises the full detection pipeline"),
    ]
    col_x = [0.3, 1.9, 3.7, 7.0]
    col_w = [1.55, 1.75, 3.2, 2.7]
    table(slide, headers, rows, col_x, col_w, header_y=1.0, row_h=1.4, start_y=1.4, font_size=9.5)
    textbox(slide, 0.3, 4.5, 9.2, 0.6,
            "Both paths feed the same downstream schema (procedures, message_log) so all 6 PCAP detectors "
            "work identically regardless of source.",
            10.5, GREY, wrap=True)


# ── Slide 12: Q&A Cheat Sheet ─────────────────────────────────────────────

def s_qa(prs):
    slide = new_slide(prs)
    slide_title(slide, "Viva Cheat-Sheet — Likely Questions & One-Line Answers")
    qa = [
        ("Why not just use one ML model for everything?",
         "Each detector has a blind spot — IF misses local/density anomalies, OC-SVM misses sequence order, "
         "rules miss novel patterns. The ensemble's overlap is the point: agreement = high confidence."),
        ("Why unsupervised detectors — why not train a classifier?",
         "No labeled 'anomaly' dataset exists for telecom captures. Unsupervised methods need zero labels and "
         "still produce a ranked severity via scores/thresholds."),
        ("Why Isolation Forest over deep learning for PCAP features?",
         "Procedure counts per capture are small (often <20) — deep models would overfit. IF is O(n log n), "
         "no distributional assumption, robust on small N."),
        ("Why CUSUM AND Bollinger AND Trend — aren't they redundant?",
         "Different time horizons: Bollinger = short-term spikes, CUSUM = sustained drift, Trend = long-term "
         "slope. Each is the textbook tool for its horizon."),
        ("Why Prophet over ARIMA?",
         "ARIMA needs manual stationarity tests (ADF) and ACF/PACF tuning per cell x KPI combination — not "
         "scalable. Prophet auto-tunes changepoints and gives uncertainty bounds out of the box."),
        ("Why a local phi3:mini LLM instead of GPT-4/Claude API?",
         "Offline requirement, CPU-only feasibility, zero API cost for batch runs, and data privacy for "
         "operator captures. RAG retrieval supplies the facts; the LLM just synthesizes."),
        ("What happens if Ollama isn't installed / fails?",
         "Deterministic rule-based fallback builds the same structured explanation from the retrieved 3GPP "
         "chunks — the system degrades gracefully, never crashes."),
        ("Why FAISS instead of a hosted vector DB (Pinecone/Chroma)?",
         "Only ~36 knowledge chunks — exact search (IndexFlatIP) is instant and needs no server, no network, "
         "fully reproducible offline."),
        ("How does the system 'learn' over time without labeled data?",
         "Engineer feedback (true/false positive ratings) feeds a proportional controller that nudges detector "
         "thresholds/sensitivity — not full retraining, since there's no ground-truth loss to retrain against."),
        ("Why rule-based correlation instead of a graph/clustering model?",
         "With 3 sources and a shared key (cell_id), a learned model is overkill and unexplainable. Bucketing "
         "mirrors how a NOC engineer manually cross-checks alarms."),
    ]
    col_x = [0.3, 4.6]
    col_w = [4.0, 5.3]
    for i, (q, a) in enumerate(qa):
        top = 0.78 + i * 0.475
        textbox(slide, col_x[0], top, col_w[0], 0.47, "Q: " + q, 8.7, ACCENT, bold=True, wrap=True)
        textbox(slide, col_x[1], top, col_w[1], 0.47, "A: " + a, 8.7, WHITE, wrap=True)


# ── Slide 13: Closing — design philosophy ─────────────────────────────────

def s_closing(prs):
    slide = new_slide(prs)
    slide_title(slide, "Design Philosophy — One Sentence Per Layer")
    bullets(slide, 0.4, 0.9, 9.2, 4.3, [
        "Detection: complementary unsupervised ensembles (tabular + sequential + rule-based) — no labels, "
        "agreement raises confidence, disagreement is itself informative.",
        "Correlation: deterministic rule-based bucketing — auditable, mirrors manual NOC practice.",
        "Prediction: two independent forecasters (decomposition + neural) cross-validate each other.",
        "Explanation: small local LLM + RAG, with a rule-based fallback — offline, private, never fails closed.",
        "MLOps: feedback-driven threshold tuning (P-controller) — explainable, stable, no retraining infra needed.",
        "Every choice optimizes for the same constraints: CPU-only, offline, no labeled data, fully explainable "
        "to a reviewer.",
    ], 13, WHITE, line_h=0.62)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    s_title(prs)
    s_overview(prs)
    s_pcap_ensemble(prs)
    s_kpi_ensemble(prs)
    s_stats_ensemble(prs)
    s_prediction(prs)
    s_correlation(prs)
    s_rag(prs)
    s_llm(prs)
    s_retrainer(prs)
    s_parsing(prs)
    s_qa(prs)
    s_closing(prs)

    out = Path("docs/Algorithm_Justification.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
