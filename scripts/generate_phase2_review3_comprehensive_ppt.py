"""
Generate the Phase II — Third Review COMPREHENSIVE PPT — Unified Telecom Analyzer.

Flow:
  Part A  — Phase II First Review  recap  (novelty proposal, 40% code)
  Part B  — Phase II Second Review recap  (modified algo, 80% code, intermediate results)
  Part C  — Phase II Third Review  (full checklist):
              Title & Abstract | Overall Design (Ph I + Ph II) | Experimental Results
              Performance Evaluation | Comparison with Existing System | References
              100% Code Implementation – Demo | Journal Paper Publication Proof
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)   # cyan
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2D, 0xC6, 0x53)
ORANGE  = RGBColor(0xFF, 0x94, 0x00)
RED     = RGBColor(0xFF, 0x4D, 0x4D)
GREY    = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE  = RGBColor(0xA0, 0x60, 0xFF)
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)
PINK    = RGBColor(0xFF, 0x4D, 0x9F)
TEAL    = RGBColor(0x2E, 0xC4, 0xB6)

# Section banner colours
PART_A_COLOR = RGBColor(0xFF, 0x94, 0x00)   # orange  — Review 1
PART_B_COLOR = RGBColor(0xA0, 0x60, 0xFF)   # purple  — Review 2
PART_C_COLOR = RGBColor(0x00, 0xB4, 0xD8)   # cyan    — Review 3


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size, color, bold_first=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if bold_first and i == 0:
            r.font.bold = True
    return tb


def table_header(slide, col_x, col_w, headers, top, size=10, color=ACCENT):
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], top, col_w[i], 0.3, h, size, color, bold=True, wrap=True)


def table_rows(slide, col_x, col_w, rows, top0, row_h=0.42, size=9.5,
               default_color=WHITE, last_col_green=False):
    for r, row in enumerate(rows):
        ty = top0 + r * row_h
        for i, val in enumerate(row):
            c = GREEN if (last_col_green and i == len(row) - 1 and val in ("Done", "✓")) else default_color
            textbox(slide, col_x[i], ty, col_w[i], row_h - 0.04, val, size, c, wrap=True)


def slide_title(slide, text, color=ACCENT):
    textbox(slide, 0.3, 0.12, 9.4, 0.48, text, 21, color, bold=True)


def section_badge(slide, label, color):
    """Small coloured pill in top-right corner to mark the section."""
    textbox(slide, 8.0, 0.12, 1.8, 0.36, label, 8.5, color, bold=True, align=PP_ALIGN.RIGHT)


def hline(slide, y, color=ACCENT):
    line = slide.shapes.add_shape(1, Inches(0.3), Inches(y), Inches(9.4), Emu(14000))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def footer(slide, n, label="Unified Telecom Analyzer — Phase II Third Review (Comprehensive)"):
    textbox(slide, 0.3, 5.28, 7.5, 0.28, label, 7.5, GREY)
    textbox(slide, 9.3, 5.28, 0.5, 0.28, str(n), 7.5, GREY, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  COVER SLIDE
# ══════════════════════════════════════════════════════════════════════════════

def s_cover(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 0.55, 9.0, 0.8,
            "Unified Telecom Analyzer", 38, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.45, 9.0, 0.45,
            "Multi-Modal Anomaly Detection & Explanation Framework for 5G Networks",
            16, WHITE, align=PP_ALIGN.CENTER)
    hline(slide, 2.0, ACCENT)
    textbox(slide, 0.5, 2.08, 9.0, 0.38,
            "PHASE II — THIRD REVIEW  (Final Submission)",
            17, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.52, 9.0, 0.38,
            "End-to-End Presentation: Review 1  →  Review 2  →  Review 3",
            13, ORANGE, align=PP_ALIGN.CENTER)
    hline(slide, 3.0, GREY)

    # Three columns for the three reviews
    cols = [
        (0.6,  "Phase II\nReview 1",  PART_A_COLOR, "Novelty Proposal\n40% Code"),
        (3.8,  "Phase II\nReview 2",  PART_B_COLOR, "Modified Algo\n80% Code"),
        (7.0,  "Phase II\nReview 3",  PART_C_COLOR, "Full Results\n100% Code + Journal"),
    ]
    for x, label, clr, sub in cols:
        textbox(slide, x, 3.1, 2.8, 0.5, label, 14, clr, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, x, 3.65, 2.8, 0.4, sub, 10, WHITE, align=PP_ALIGN.CENTER)

    textbox(slide, 0.5, 4.35, 9.0, 0.3,
            "M.Tech Data Science  |  PES University, Bangalore — Great Learning",
            11, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.68, 9.0, 0.28,
            "Candidate: Vikas  |  Repository: telecom-analyzer  |  Python · scikit-learn · "
            "PyTorch · Prophet · FAISS · Ollama · Streamlit · FastAPI",
            9.5, GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  PART A — PHASE II FIRST REVIEW
# ══════════════════════════════════════════════════════════════════════════════

def s_part_a_banner(prs):
    slide = new_slide(prs)
    set_bg(slide, RGBColor(0x0A, 0x12, 0x1E))
    textbox(slide, 0.5, 1.5, 9.0, 0.7,
            "PART A", 50, PART_A_COLOR, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.3, 9.0, 0.5,
            "Phase II — First Review", 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
    hline(slide, 3.0, PART_A_COLOR)
    textbox(slide, 0.5, 3.1, 9.0, 0.9,
            "Checklist:  Phase 1 Review  ·  Novelty Proposal  ·  Title & Abstract\n"
            "Proposed System (Phase II)  ·  Algorithms / Techniques\n"
            "Expected Outcomes  ·  References  ·  40% Code (Phase II)",
            13, GREY, align=PP_ALIGN.CENTER)


def s_a_phase1_recap(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Phase I Review — What Was Built & Validated")
    section_badge(slide, "Part A — Review 1", PART_A_COLOR)
    bullets(slide, 0.4, 0.68, 9.2, 2.6, [
        "Parsers (src/parsers/): Full 3GPP stack — NAS (TS 24.501), NGAP (TS 38.413), "
        "RRC (TS 38.331), F1AP (TS 38.473), E1AP (TS 38.463), XnAP (TS 38.423); "
        "pyshark for real captures, raw discriminator fallback for synthetic.",
        "Detection — 6 PCAP detectors: Isolation Forest · Statistical/Cascade · "
        "One-Class SVM · LOF · Elliptic Envelope · LSTM Autoencoder.",
        "Detection — 6 KPI detectors: Threshold · Peer-Comparison · Trend-Regression · "
        "IQR · CUSUM · Bollinger Bands.",
        "Cross-source Event Correlation: EventRouter keyed on (cell, category, 1-h window).",
        "RAG / LLM Explanation: FAISS + all-MiniLM-L6-v2 + Ollama phi3:mini narratives.",
        "MLOps Feedback Loop: FeedbackStore + nightly_retrain.py adjusts thresholds.",
        "REST API (FastAPI): /analyze/pcap, /analyze/kpi, /health endpoints.",
        "Streamlit Dashboard: PCAP & KPI paths, 6-detector anomaly UI with rationale expander.",
        "Phase I code 100% complete; Phase I review passed before Phase II began.",
    ], 10.5, WHITE)
    footer(slide, n)


def s_a_novelty(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Novelty Proposal — Two Phase II Research Directions")
    section_badge(slide, "Part A — Review 1", PART_A_COLOR)
    textbox(slide, 0.3, 0.68, 9.4, 0.3,
            "Phase II focus: NOVELTY — improve the model and validate at real-world scale.",
            10.5, ORANGE, bold=True, wrap=True)

    textbox(slide, 0.3, 1.06, 9.2, 0.28,
            "Novelty 1: Genuine Multi-Method Forecasting (not a claim — a benchmark)", 12, GREEN, bold=True)
    bullets(slide, 0.4, 1.38, 9.2, 1.1, [
        "Phase I forecasted KPIs with Prophet + LSTM but never quantified accuracy — "
        "common in industry tools, but not publishable.",
        "Proposed: rolling-origin backtest + anomaly lead-time evaluation across Prophet, "
        "Holt-Winters (damped trend), and LSTM; report MAE / RMSE / MAPE per method.",
        "Research question: which forecasting family detects emerging degradation earliest?",
    ], 10, WHITE)

    hline(slide, 2.55, PURPLE)

    textbox(slide, 0.3, 2.65, 9.2, 0.28,
            "Novelty 2: LLM Root-Cause Agent — WHY, not just WHAT", 12, PURPLE, bold=True)
    bullets(slide, 0.4, 2.95, 9.2, 1.1, [
        "Phase I labelled anomalies per-source with no causal reasoning across sources.",
        "Proposed: a bounded ReAct agent (max 4 tool steps) backed by a 9-layer 5G "
        "protocol-stack causal rule engine + FAISS-backed 3GPP RAG; deterministic "
        "fallback so demo works without a live LLM.",
        "Research question: can an LLM agent reliably distinguish root cause from "
        "downstream symptom when correlated cross-domain anomaly groups are given?",
    ], 10, WHITE)

    hline(slide, 4.1, ACCENT)
    textbox(slide, 0.3, 4.18, 9.2, 0.3,
            "State at First Review: Both novelties ~40% implemented — "
            "LSTM + Prophet forecasters done, first ReAct loop coded, rolling-origin harness drafted.",
            10, GREY, italic=True, wrap=True)
    footer(slide, n)


def s_a_system_ph2(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Proposed System — Phase II Architecture")
    section_badge(slide, "Part A — Review 1", PART_A_COLOR)

    layers = [
        ("Phase I (unchanged)",   WHITE,
         "18-detector ensemble · EventRouter · RAG/LLM explanation · MLOps · REST API · Dashboard"),
        ("Novelty 1: Forecast Eval",  GREEN,
         "Prophet + Holt-Winters (new) + LSTM → rolling-origin backtest (MAE/RMSE/MAPE) "
         "+ evaluate_anomaly_lead_time() → anomaly lead-time table"),
        ("Novelty 2: ReAct Agent",   PURPLE,
         "9-layer causal_rules.py (PHY→MAC→RLC→PDCP→RRC→NAS→NGAP→Core→KPI) → "
         "bounded ReAct loop (4 steps, 4 tools) → 3GPP-cited root-cause narrative "
         "→ rule-based deterministic fallback (no-LLM path)"),
        ("New REST Endpoints",   ORANGE,
         "/analyze/stats (Stats domain parse+detect) · /agent/root-cause (async, "
         "cross-source correlated group → agent narrative)"),
        ("Dashboard Extension",   ACCENT,
         "Shared EventRouter per session (not per upload) · Root Cause Agent panel "
         "· simplified end-user mode · predicted anomalies fed into the router"),
    ]
    for i, (label, lcolor, desc) in enumerate(layers):
        y = 0.72 + i * 0.87
        textbox(slide, 0.3, y, 2.3, 0.78, label, 10, lcolor, bold=True, wrap=True)
        textbox(slide, 2.7, y, 7.0, 0.78, desc, 10, WHITE, wrap=True)

    footer(slide, n)


def s_a_algos(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Algorithms & Techniques — Phase II Additions")
    section_badge(slide, "Part A — Review 1", PART_A_COLOR)

    textbox(slide, 0.3, 0.65, 4.6, 0.28, "Forecasting (Novelty 1)", 11.5, GREEN, bold=True)
    bullets(slide, 0.3, 0.95, 4.6, 2.1, [
        "Prophet (Taylor & Letham 2018): additive model, "
        "seasonal decomposition, uncertainty intervals.",
        "Holt-Winters (damped trend): exponential smoothing "
        "with multiplicative seasonality — low-resource baseline for comparison.",
        "LSTM Multi-horizon: encoder produces hidden state; "
        "4 time-step direct forecast — captures sequence patterns invisible to Prophet.",
        "Evaluation: rolling_origin_backtest() — sliding N origins, "
        "training prefix grows, test window fixed at 4 h; scores MAE, RMSE, MAPE.",
        "evaluate_anomaly_lead_time(): for each ground-truth anomaly window, walks "
        "backward in steps to find earliest forecast alarm — reports lead-time hours & false-alarm count.",
    ], 9.7, WHITE)

    textbox(slide, 5.1, 0.65, 4.6, 0.28, "Root-Cause Agent (Novelty 2)", 11.5, PURPLE, bold=True)
    bullets(slide, 5.1, 0.95, 4.6, 2.1, [
        "ReAct loop (Yao et al. ICLR 2023): Reason → Act → Observe, "
        "max 4 iterations, bounded to prevent infinite loops.",
        "Tools: retrieve_3gpp_spec (FAISS RAG), query_event_history, "
        "score_causal_hypothesis, get_protocol_layer_context.",
        "9-layer causal rule engine: keyword + layer-rank matching "
        "across PHY · MAC · RLC · PDCP · RRC · NAS · NGAP · Core · KPI.",
        "RAG: FAISS index over 3GPP excerpts (TS 38.913, 38.413, 24.501, "
        "38.321, 38.331) with all-MiniLM-L6-v2 embeddings.",
        "Guardrail: any unparsable model step triggers "
        "rule_based_root_cause() — deterministic ordering by time + protocol depth.",
    ], 9.7, WHITE)

    hline(slide, 3.1, ACCENT)
    textbox(slide, 0.3, 3.2, 9.4, 0.28,
            "Expected Outcomes (stated at First Review):",
            10.5, ORANGE, bold=True)
    bullets(slide, 0.4, 3.52, 9.2, 1.4, [
        "A reproducible forecast benchmark with honest negative findings "
        "(abrupt anomalies are not forecastable with historical data alone).",
        "A live-demoable LLM agent that identifies causal root causes "
        "(not symptom labels) from correlated cross-source anomaly clusters.",
        "Second Review target: 80% code, intermediate benchmark results, "
        "3 forecasting methods scored, ReAct agent end-to-end wired.",
    ], 10, WHITE)
    footer(slide, n)


def s_a_code_status(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Phase II Code Status — First Review (40% Target)")
    section_badge(slide, "Part A — Review 1", PART_A_COLOR)

    headers = ["Component", "Status @ First Review", "Note"]
    rows = [
        ("Prophet + LSTM point forecasters",              "Done (carried from Phase I)", "~30% of forecast eval"),
        ("Holt-Winters 3rd method",                       "Proposed — not yet coded",   "Planned for Review 2"),
        ("Rolling-origin backtest harness",               "Draft — not yet tested",      "forecast_eval.py skeleton"),
        ("anomaly_lead_time evaluator",                   "Draft — single manual pass",  "needs formal loop"),
        ("9-layer causal_rules.py",                       "~60% done",                   "layer matching needs tuning"),
        ("Bounded ReAct loop (4-step, 4 tools)",          "First loop coded",            "parse guardrail pending"),
        ("rule_based_root_cause() fallback",              "Not yet implemented",         "Critical for demo reliability"),
        ("/analyze/stats + /agent/root-cause endpoints",  "Not yet wired",               "needs async threading"),
        ("Dashboard shared session router",               "Not yet — per-upload state",  "Review 2 target"),
        ("Automated tests for agent / forecast eval",     "Not yet",                     "24 tests planned Review 2"),
    ]
    col_x = [0.3, 4.5, 6.7]
    col_w = [4.1, 2.1, 2.75]
    table_header(slide, col_x, col_w, headers, 0.68)
    for r, row in enumerate(rows):
        ty = 1.02 + r * 0.4
        textbox(slide, col_x[0], ty, col_w[0], 0.37, row[0], 9, WHITE, wrap=True)
        clr = GREEN if "Done" in row[1] else (ORANGE if "Draft" in row[1] or "%" in row[1] else RED)
        textbox(slide, col_x[1], ty, col_w[1], 0.37, row[1], 9, clr, wrap=True)
        textbox(slide, col_x[2], ty, col_w[2], 0.37, row[2], 8.5, GREY, wrap=True)

    textbox(slide, 0.3, 5.1, 9.4, 0.28,
            "Overall: ~40% of planned Phase II scope complete at First Review.",
            10, ORANGE, bold=True, italic=True, wrap=True)
    footer(slide, n)


# ══════════════════════════════════════════════════════════════════════════════
#  PART B — PHASE II SECOND REVIEW
# ══════════════════════════════════════════════════════════════════════════════

def s_part_b_banner(prs):
    slide = new_slide(prs)
    set_bg(slide, RGBColor(0x0A, 0x12, 0x1E))
    textbox(slide, 0.5, 1.5, 9.0, 0.7,
            "PART B", 50, PART_B_COLOR, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.3, 9.0, 0.5,
            "Phase II — Second Review", 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
    hline(slide, 3.0, PART_B_COLOR)
    textbox(slide, 0.5, 3.1, 9.0, 0.9,
            "Checklist:  Title & Abstract  ·  Modified Algorithm Design\n"
            "Contribution of the Candidate  ·  Results Obtained (Intermediate)\n"
            "References  ·  80% of Code (Phase II)",
            13, GREY, align=PP_ALIGN.CENTER)


def s_b_abstract(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Phase II Second Review — Title & Abstract")
    section_badge(slide, "Part B — Review 2", PART_B_COLOR)
    textbox(slide, 0.3, 0.7, 9.4, 0.36,
            "Title: Genuine Forecast Evaluation and a Guardrailed LLM Root-Cause Agent "
            "for a Multi-Domain 5G Anomaly Detection Framework",
            12, GREEN, bold=True, wrap=True)
    text = (
        "First Review identified two gaps in Phase I: forecasting was claimed but never "
        "quantified, and anomalies were detected in isolation with no causal explanation "
        "across sources. This review reports both gaps closed.\n\n"
        "A rolling-origin backtest and an anomaly lead-time evaluator now score three "
        "forecasting methods (Prophet, Holt-Winters, direct multi-horizon LSTM) on "
        "MAE/RMSE/MAPE and on whether they flag injected anomalies before onset — "
        "including an honest negative finding: abrupt step-change anomalies have no "
        "learnable precursor, contrasted with a positive demonstration on a genuine "
        "leading-trend series.\n\n"
        "A bounded ReAct agent, backed by a 9-layer 5G causal-rule engine and RAG-grounded "
        "3GPP citations, explains WHY a correlated group of cross-source anomalies occurred "
        "— with a deterministic rule-based fallback so the system is demoable with or "
        "without a live LLM. Both novelties are wired end-to-end: two new REST endpoints "
        "(/analyze/stats, /agent/root-cause) and a shared cross-upload session router "
        "in the dashboard."
    )
    textbox(slide, 0.3, 1.18, 9.4, 3.95, text, 11.5, WHITE, wrap=True)
    footer(slide, n)


def s_b_modified_algo(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Modified Algorithm Design — Changes Since First Review")
    section_badge(slide, "Part B — Review 2", PART_B_COLOR)
    textbox(slide, 0.3, 0.65, 9.4, 0.26,
            "Structural changes per component:", 10.5, WHITE, italic=True, wrap=True)

    headers = ["Component", "First Review", "Second Review (this deck)"]
    rows = [
        ("Forecasting methods",
         "Prophet + LSTM, ad-hoc accuracy checks",
         "Added Holt-Winters (damped trend) as 3rd method; formal rolling_origin_backtest() "
         "sliding N origins per (method, cell, column) → MAE/RMSE/MAPE"),
        ("Lead-time evaluation",
         "Single manual pass over 3 injected anomalies",
         "evaluate_anomaly_lead_time() walks backward from each ground-truth window "
         "in fixed steps + scans for false alarms outside any window"),
        ("Root-cause reasoning",
         "Single-shot LLM prompt, no tool use, no fallback",
         "Bounded ReAct loop (max 4 steps, 4 typed tools); parse failure falls through "
         "to fixed deterministic tool order — never aborts the demo"),
        ("Causal knowledge base",
         "Implicit in the prompt string",
         "Explicit 9-layer causal_rules.py (PHY→KPI) with keyword + layer-rank matching"),
        ("No-LLM fallback",
         "None — agent required Ollama",
         "rule_based_root_cause(): deterministic chain ordering by time + protocol depth, "
         "real RAG retrieval per hop — 3GPP citation guaranteed without a live model"),
        ("System wiring",
         "Agent + forecaster importable; one dashboard panel",
         "/analyze/stats + /agent/root-cause (async via asyncio.to_thread); dashboard "
         "keeps one EventRouter per session across uploads"),
    ]
    col_x = [0.3, 2.1, 5.55]
    col_w = [1.7, 3.35, 4.05]
    table_header(slide, col_x, col_w, headers, 0.97, size=10)
    for r, row in enumerate(rows):
        ty = 1.3 + r * 0.675
        for i, val in enumerate(row):
            textbox(slide, col_x[i], ty, col_w[i], 0.64, val, 8.7, WHITE, wrap=True)
    footer(slide, n)


def s_b_contribution(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Contribution of the Candidate — Second Review")
    section_badge(slide, "Part B — Review 2", PART_B_COLOR)
    bullets(slide, 0.4, 0.7, 9.2, 4.5, [
        "Built src/detection/forecast_eval.py from scratch: rolling-origin backtest + "
        "anomaly lead-time evaluator + run_benchmark_report() single entry point, plus "
        "the ground-truth event table for the 10-hour KPI sample.",
        "Designed and implemented src/agent/ (schemas.py, tools.py, causal_rules.py, "
        "rule_based.py, react_agent.py): bounded ReAct agent with parse-failure guardrail, "
        "4 individually-testable tool wrappers, 9-layer causal-rule engine, deterministic "
        "RAG-grounded fallback — adapted from earlier reverted RCA attempt to EventRouter's "
        "current event schema.",
        "Added Holt-Winters as 3rd forecasting method; refactored threshold/severity "
        "helpers in predictor.py so forecast_eval.py, the pipeline, and the dashboard "
        "all share one implementation.",
        "Extended REST API (src/api/main.py): new /analyze/stats endpoint and "
        "/agent/root-cause that combines job results into a shared EventRouter and "
        "runs the agent over cross-source correlated groups, off the event loop.",
        "Reworked dashboard session state to one shared EventRouter across uploads; "
        "added Root Cause Agent panel and simplified end-user mode; routed predicted "
        "anomalies into the shared router.",
        "Wrote 24 new tests (test_root_cause_agent.py, test_forecast_eval.py, "
        "test_api_root_cause.py, test_pipeline_prediction.py); extended existing suites; "
        "found and fixed 2 pre-existing bugs (PDF-export crash, freq_s divide-by-zero).",
    ], 10.8, WHITE)
    footer(slide, n)


def s_b_results(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Results Obtained — Intermediate (Second Review)")
    section_badge(slide, "Part B — Review 2", PART_B_COLOR)

    headers = ["Metric", "Value"]
    rows = [
        ("Full automated test suite",            "83 passed, 1 skipped (needs live Ollama), 0 failed  [→ Review 3 final: 173 passed]"),
        ("New tests added this review",           "24  (agent 13, forecast_eval 6, API root-cause 3, pipeline 2)"),
        ("Regressions vs Phase I / Review 1",     "0"),
        ("Forecasting methods benchmarked",       "3 — Prophet, Holt-Winters, direct multi-horizon LSTM"),
        ("Forecast scoring",                      "Rolling-origin MAE / RMSE / MAPE per (method, cell, column)"),
        ("Anomaly lead-time result",              "Abrupt step-change: undetectable ahead of onset (honest finding). "
                                                  "Gradual trend: correctly flagged with real lead-time — positive control."),
        ("Live cross-source demo",                "PCI_5 KPI + Stats: dl_bler (PHY) → Handover SR + Cell Availability "
                                                  "drops correctly attributed (TS 38.913, TS 38.413 citations)."),
    ]
    col_x = [0.3, 3.6]
    col_w = [3.2, 5.9]
    table_header(slide, col_x, col_w, headers, 0.68, size=11)
    for r, row in enumerate(rows):
        ty = 1.02 + r * 0.57
        textbox(slide, col_x[0], ty, col_w[0], 0.53, row[0], 10, WHITE, wrap=True)
        textbox(slide, col_x[1], ty, col_w[1], 0.53, row[1], 10, WHITE, wrap=True)

    footer(slide, n)


def s_b_code_status(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Phase II Code Completion — Second Review (80% Target)")
    section_badge(slide, "Part B — Review 2", PART_B_COLOR)
    headers = ["Item", "Status"]
    rows = [
        ("Prophet + Holt-Winters + LSTM point forecasters",                          "Done"),
        ("Rolling-origin backtest + anomaly lead-time evaluation harness",            "Done"),
        ("9-layer causal rule engine (src/agent/causal_rules.py)",                   "Done"),
        ("Bounded ReAct root-cause agent + parse-failure guardrail",                 "Done"),
        ("Deterministic RAG-grounded rule-based fallback (no-LLM path)",             "Done"),
        ("/analyze/stats + /agent/root-cause REST endpoints",                        "Done"),
        ("Dashboard: shared session router, Root Cause Agent panel, simple mode",    "Done"),
        ("24 new automated tests, 0 regressions (83 passed at Review 2 → 173 by Review 3)", "Done"),
        ("Benchmark vs published single-method baselines",                            "Pending — Third Review"),
        ("Performance / scalability profiling at real-world scale (64 UE / 6 hr)",   "Pending — Third Review"),
        ("Journal paper draft (Phase I + II)",                                        "Pending — Third Review"),
    ]
    col_x = [0.3, 7.3]
    col_w = [6.85, 2.35]
    table_header(slide, col_x, col_w, headers, 0.62, size=10.5)
    for r, row in enumerate(rows):
        ty = 0.96 + r * 0.38
        textbox(slide, col_x[0], ty, col_w[0], 0.35, row[0], 9.2, WHITE, wrap=True)
        clr = GREEN if row[1] == "Done" else ORANGE
        textbox(slide, col_x[1], ty, col_w[1], 0.35, row[1], 9.2, clr, bold=True)

    textbox(slide, 0.3, 5.1, 9.4, 0.28,
            "8 of 11 items done  →  ~80% of planned Phase II scope complete; "
            "remaining items are Third Review deliverables.",
            9.5, GREY, italic=True, wrap=True)
    footer(slide, n)


# ══════════════════════════════════════════════════════════════════════════════
#  PART C — PHASE II THIRD REVIEW (FINAL)
# ══════════════════════════════════════════════════════════════════════════════

def s_part_c_banner(prs):
    slide = new_slide(prs)
    set_bg(slide, RGBColor(0x0A, 0x12, 0x1E))
    textbox(slide, 0.5, 1.5, 9.0, 0.7,
            "PART C", 50, PART_C_COLOR, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.3, 9.0, 0.5,
            "Phase II — Third Review (Final)", 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
    hline(slide, 3.0, PART_C_COLOR)
    textbox(slide, 0.5, 3.1, 9.0, 1.1,
            "Checklist:  Title & Abstract  ·  Overall Design (Phase I + Phase II)\n"
            "Experimental Results  ·  Performance Evaluation\n"
            "Comparison with Existing System  ·  References\n"
            "100% Code Implementation – Demo  ·  Journal Paper Publication Proof",
            13, GREY, align=PP_ALIGN.CENTER)


def s_c_title_abstract(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Title & Abstract — Final")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)
    textbox(slide, 0.3, 0.68, 9.4, 0.36,
            "Unified Telecom Analyzer: An 18-Detector Multi-Domain Framework for Anomaly Detection, "
            "Cross-Domain Correlation, and LLM-Aided Explanation in 5G Networks",
            12.5, GREEN, bold=True, wrap=True)
    text = (
        "Modern 5G networks generate three concurrent data streams — PCAP protocol traces, "
        "per-cell KPI exports, and gNB L1/L2 statistics — yet existing monitoring tools "
        "treat each source in isolation, missing cross-domain fault signatures. This work "
        "presents a unified, open-source pipeline that parses all three domains against the "
        "full 3GPP specification stack, runs 18 detectors across PCAP, KPI, and Stats domains "
        "(6 methods per domain) covering statistical, density, kernel, and neural families, "
        "routes anomalies through a shared event correlator keyed on (cell, category, time-window), "
        "forecasts near-term degradation with Holt-Winters and LSTM (benchmarked; Prophet "
        "implemented but not executed in the authoritative benchmark environment), and provides "
        "natural-language root-cause explanations via FAISS-backed RAG → Ollama LLM "
        "with a deterministic rule-based fallback.\n\n"
        "Phase II was evaluated on a controlled synthetic testbed: 4 gNBs, "
        "16 cells, 64 UEs, 6-hour diurnal run with 3 injected fault scenarios (congestion, "
        "outage, drift). Under the final cell-level evaluation (3 fault cells / 16 total), "
        "Stats Count-Threshold achieved Precision=1.00, Recall=1.00, F1=1.00. "
        "The Full Ensemble achieved Recall=1.00 (all faults detected) with Precision=0.1875, "
        "F1=0.316. The system correlates anomaly events across domains, provides explainable "
        "RCA narratives, and completes end-to-end in 45.6 s on a CPU-only host. "
        "173 automated tests, 0 failures."
    )
    textbox(slide, 0.3, 1.18, 9.4, 3.95, text, 11.2, WHITE, wrap=True)
    footer(slide, n)


def s_c_overall_design(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Overall System Design — Phase I + Phase II")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)

    layers = [
        ("INGEST",    ACCENT,
         "PCAP Upload  |  KPI Excel/CSV  |  Stats Parquet/CSV — via Streamlit UI, FastAPI REST, or CLI"),
        ("PARSE",     PURPLE,
         "pcap_parser_real.py (pyshark): NAS·NGAP·RRC·F1AP·E1AP·XnAP vs full 3GPP stack  |  "
         "kpi_parser.py (vs 3GPP KPI catalogue)  |  stats_parser.py (srsRAN/OAI/NIST/Parquet)"),
        ("DETECT",    ORANGE,
         "18 detectors across 3 domains — PCAP: Isolation Forest · Statistical/Cascade · OC-SVM · LOF · "
         "Elliptic Env. · LSTM-AE  |  KPI: Threshold · Peer-Comparison · Trend-Regression · IQR · CUSUM · "
         "Bollinger  |  Stats: same family, adapted for L1/L2 DU/CU metrics"),
        ("CORRELATE", YELLOW,
         "EventRouter: (cell, category) key, 1-h sliding window → cross-domain anomaly clusters  |  "
         "1,427 anomaly events (KPI: 1,316 + Stats: 111) — 100% cross-source cluster assignment at 64-UE scale"),
        ("PREDICT",   GREEN,
         "Holt-Winters (damped trend) + LSTM (sequence) — 4-hr ahead, rolling-origin backtest + "
         "lead-time evaluator  |  HW MAE=15.54, LSTM MAE=20.24 (seed=42)  |  "
         "Prophet implemented — not executed in authoritative benchmark environment"),
        ("EXPLAIN",   TEAL,
         "FAISS + all-MiniLM-L6-v2 RAG → Ollama phi3:mini ReAct agent → natural-language root-cause  |  "
         "rule_based_root_cause() deterministic fallback — demoable without live LLM"),
        ("MLOPS",     PINK,
         "FeedbackStore → nightly_retrain.py adjusts detector thresholds from false-positive feedback"),
        ("INTERFACES", GREY,
         "Streamlit Dashboard  |  FastAPI REST (POST /analyze/pcap · /analyze/kpi · /analyze/stats · "
         "POST /agent/root-cause · GET /health)  |  CLI (python -m src.orchestrator.pipeline)"),
    ]

    for i, (label, lcolor, desc) in enumerate(layers):
        y = 0.65 + i * 0.605
        textbox(slide, 0.3, y, 1.65, 0.54, label, 9.5, lcolor, bold=True, wrap=True)
        textbox(slide, 2.05, y, 7.6, 0.54, desc, 8.8, WHITE, wrap=True)

    footer(slide, n)


def s_c_experimental_results(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Experimental Results")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)
    textbox(slide, 0.3, 0.62, 9.4, 0.26,
            "Dataset: 4 gNB / 16 cell / 64 UE / 6-hr run — 3 concurrent injected faults: "
            "congestion (PCI_3), outage (PCI_12), drift (PCI_8).  Pipeline: --no-llm flag.",
            9.5, GREY, wrap=True)

    headers = ["Domain", "Input size", "Anomaly records", "Top detectors"]
    rows = [
        ("PCAP",  "901 pkts / 901 events",    "16",    "Statistical/Cascade (9) · LSTM-AE (51 frame-level) · IF (3)"),
        ("KPI",   "5,760 rows × 26 cols",     "1,316", "Threshold (833) · CUSUM (212) · Bollinger (196) · IQR (117)"),
        ("Stats", "5,760 rows × 15 cols",     "111",   "Trend-Regression (88) · Peer-Comparison (15) · IQR (8)"),
    ]
    col_x = [0.3, 2.05, 3.8, 5.5]
    col_w = [1.65, 1.65, 1.65, 4.1]
    table_header(slide, col_x, col_w, headers, 0.94, size=10)
    for r, row in enumerate(rows):
        ty = 1.27 + r * 0.55
        for i, val in enumerate(row):
            textbox(slide, col_x[i], ty, col_w[i], 0.51, val, 9.5, WHITE, wrap=True)

    textbox(slide, 0.3, 2.98, 9.2, 0.28,
            "Cross-Domain Correlation (shared EventRouter — PCAP + KPI + Stats):",
            11, ACCENT, bold=True)
    bullets(slide, 0.4, 3.29, 9.0, 1.25, [
        "EventRouter processed 1,427 anomaly events (KPI: 1,316 + Stats: 111). "
        "Event-level cross-source assignment rate: 100% (1,427/1,427). "
        "Cell-level: all 16 cells have ≥2 source domains — note this includes 13 normal "
        "cells where KPI fires due to FDD/TDD capacity calibration differences.",
        "All 3 injected fault cells (PCI_3, PCI_12, PCI_8) are correctly surfaced: "
        "congestion (CUSUM + Bollinger), outage (Threshold + Peer-Comparison), drift (Trend-Regression).",
        "Cross-source coverage alone does not imply fault attribution — "
        "the Stats Count-Threshold step (F1=1.00) is needed to isolate the 3 true fault cells "
        "from the 13 normal cells that also generate cross-source events.",
    ], 10, WHITE)
    footer(slide, n)


def s_c_performance(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Performance Evaluation")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)

    # ── Detection benchmark header ──────────────────────────────────────────
    textbox(slide, 0.3, 0.62, 9.4, 0.24,
            "Final Cell-Level Detection Benchmark  "
            "(Ground truth: 3 injected fault cells / 16 total cells, seed=42):",
            10.5, ACCENT, bold=True, wrap=True)

    headers = ["Method", "TP", "FP", "FN", "Precision", "Recall", "F1"]
    rows = [
        ("Threshold-only (KPI)",         "3", "4",  "0", "0.4286", "1.000", "0.600"),
        ("IQR-only (KPI)",               "3", "13", "0", "0.1875", "1.000", "0.316"),
        ("Full Ensemble (KPI+Stats)",    "3", "13", "0", "0.1875", "1.000", "0.316"),
        ("Ensemble + EventRouter",       "3", "13", "0", "0.1875", "1.000", "0.316"),
        ("Stats Count-Threshold ★",     "3", "0",  "0", "1.000",  "1.000", "1.000"),
    ]
    col_x = [0.3, 3.25, 3.85, 4.45, 5.1, 6.2, 7.3]
    col_w = [2.85, 0.5,  0.5,  0.55, 1.0, 1.0, 1.0]
    table_header(slide, col_x, col_w, headers, 0.90, size=9.5)
    row_colors = [WHITE, WHITE, WHITE, WHITE, GREEN]
    for r, row in enumerate(rows):
        ty = 1.20 + r * 0.37
        for i, val in enumerate(row):
            textbox(slide, col_x[i], ty, col_w[i], 0.34, val,
                    9.0, row_colors[r], bold=(r == 4), wrap=True)

    textbox(slide, 0.3, 3.10, 9.4, 0.22,
            "★ Best evaluated method: Stats Count-Threshold F1=1.00  |  "
            "All methods achieve Recall=1.00 (all 3 fault cells detected)  |  "
            "KPI ensemble FP rate reflects TDD/FDD capacity calibration differences in this dataset.",
            8.5, GREY, italic=True, wrap=True)

    # ── Scalability ─────────────────────────────────────────────────────────
    textbox(slide, 0.3, 3.40, 9.4, 0.24,
            "Pipeline Scalability — 64-UE / 6-hr Run  (CPU-only host, tracemalloc memory):",
            10.5, ACCENT, bold=True, wrap=True)

    # AUTHORITATIVE VALUES from the original Python 3.13.13 benchmark run (git a516829).
    # results/scalability_results.json reflects a later Python 3.12.3 re-run (40.1 s total).
    # The PPT intentionally uses the original run to maintain consistency with the
    # submitted experiment_manifest.json. Any future re-benchmark should update BOTH.
    sc_rows = [
        ("PCAP parse (901 pkts, pyshark)",        "2.6 s",  "1.2 MB"),
        ("KPI parse (5,760 rows × 26 cols)",       "1.1 s",  "11.1 MB"),
        ("KPI detection (6 detectors)",            "32.7 s", "7.8 MB"),
        ("Stats parse + detection (5,760 rows)",   "5.0 s",  "17.0 MB"),
        ("Cross-domain correlation (1,427 events)","4.2 s",  "1.6 MB"),
        ("TOTAL end-to-end pipeline",              "45.6 s", "~21 MB peak (tracemalloc)"),
    ]
    ph = ["Stage", "Latency", "Memory (incremental)"]
    pc = [0.3, 6.1, 7.6]
    pw = [5.7, 1.4, 2.0]
    table_header(slide, pc, pw, ph, 3.68, size=9.5)
    for r, row in enumerate(sc_rows):
        ty = 3.96 + r * 0.245
        c = YELLOW if row[0].startswith("TOTAL") else WHITE
        for i, val in enumerate(row):
            textbox(slide, pc[i], ty, pw[i], 0.23, val, 8.5, c)
    footer(slide, n)


def s_c_comparison(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Comparison with Existing Systems")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)

    headers = ["Property", "Threshold-\nonly", "Isolation\nForest", "One-Class\nSVM",
               "Our Ensemble\n(6 × 3 domains)"]
    rows = [
        ("Anomaly families covered",
         "Point spikes only",
         "Point + moderate cluster",
         "Boundary outliers",
         "Spike · cluster · sequence · trend · drift — all families"),
        ("Drift / slow degradation",
         "Misses",
         "Misses",
         "Misses",
         "Caught by Trend-Regression + CUSUM (88 + 212 this run)"),
        ("Sequence/order violations",
         "None",
         "None",
         "None",
         "LSTM-AE reconstructs procedure order; 51 frame-level alerts"),
        ("Cross-domain correlation",
         "None",
         "None",
         "None",
         "EventRouter: PCAP + KPI + Stats; 1,427 events, 100% cross-source assignment"),
        ("Explainability",
         "Rule string",
         "Feature importance",
         "Support-vector dist.",
         "RAG + LLM narrative + rationale expander in dashboard"),
        ("Cell-level Precision / F1\n(final benchmark)",
         "0.4286 / 0.600",
         "Not evaluated\nin final benchmark",
         "Not evaluated\nin final benchmark",
         "0.1875 / 0.316 (Ensemble+Router)\n1.000 / 1.000 (Stats Count-Threshold)"),
        ("Multi-source domain",
         "KPI/Stats only",
         "Single source",
         "Single source",
         "PCAP + KPI + Stats — unified, one pipeline"),
        ("Open-source + demoable",
         "Partial",
         "Partial",
         "Partial",
         "Full — Streamlit + FastAPI + CLI, 173 tests passed, 0 failed"),
    ]

    col_x = [0.3, 2.65, 4.15, 5.65, 7.15]
    col_w = [2.25, 1.4, 1.4, 1.4, 2.45]
    table_header(slide, col_x, col_w, headers, 0.63, size=9)
    for r, row in enumerate(rows):
        ty = 1.02 + r * 0.545
        for i, val in enumerate(row):
            c = GREEN if i == 4 else WHITE
            textbox(slide, col_x[i], ty, col_w[i], 0.52, val, 8.5, c, wrap=True)
    textbox(slide, 0.3, 5.30, 9.4, 0.22,
            "★ Note: Isolation Forest and One-Class SVM columns show algorithmic baseline "
            "characteristics only — neither was part of the final cell-level benchmark run.",
            8.0, GREY, italic=True, wrap=True)
    footer(slide, n)


def s_c_code_demo(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "100% Code Implementation — Demo")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)

    textbox(slide, 0.3, 0.63, 9.4, 0.26,
            "All planned modules are implemented, tested, and runnable end-to-end.",
            10, GREY, italic=True, wrap=True)

    headers = ["Module / File", "~LOC", "Status"]
    rows = [
        ("src/parsers/pcap_parser_real.py — NAS/NGAP/RRC/F1AP/E1AP/XnAP (pyshark)", "620", "Done"),
        ("src/parsers/kpi_parser.py — Excel/CSV KPI vs 3GPP catalogue",               "310", "Done"),
        ("src/parsers/stats_parser.py — srsRAN/OAI/NIST/Parquet DU-CU stats",        "280", "Done"),
        ("src/detection/detectors.py — 6 PCAP detectors + ensemble merge",            "540", "Done"),
        ("src/detection/kpi_detector.py — 6 KPI/Stats detectors",                     "480", "Done"),
        ("src/detection/forecast_eval.py — rolling-origin + lead-time evaluator",     "310", "Done"),
        ("src/orchestrator/event_router.py — cross-domain EventRouter",               "210", "Done"),
        ("src/prediction/prophet_model.py + lstm_model.py — 4-hr forecasting",       "390", "Done"),
        ("src/agent/ — schemas, tools, causal_rules, rule_based, react_agent",        "520", "Done"),
        ("src/explanation/rag_engine.py — FAISS + LLM + rule fallback",              "260", "Done"),
        ("src/mlops/feedback_store.py + nightly_retrain.py — MLOps loop",            "195", "Done"),
        ("src/orchestrator/pipeline.py — CLI orchestrator (run_pipeline)",            "230", "Done"),
        ("src/api/main.py — FastAPI REST (5 endpoints)",                              "210", "Done"),
        ("src/dashboard/app.py — Streamlit (PCAP + KPI + Agent panel)",              "910", "Done"),
        ("tests/ — 173 automated tests (pytest), 0 failures, 1 skipped",              "680", "Done"),
    ]
    col_x = [0.3, 7.5, 8.3]
    col_w = [7.1, 0.7, 1.3]
    table_header(slide, col_x, col_w, headers, 0.93, size=9.5)
    for r, row in enumerate(rows):
        ty = 1.22 + r * 0.267
        textbox(slide, col_x[0], ty, col_w[0], 0.25, row[0], 8.0, WHITE, wrap=True)
        textbox(slide, col_x[1], ty, col_w[1], 0.25, row[1], 8.0, GREY)
        textbox(slide, col_x[2], ty, col_w[2], 0.25, row[2], 8.0, GREEN, bold=True)

    textbox(slide, 0.3, 5.22, 9.4, 0.26,
            "Demo:  make dev  (Streamlit)  |  make api  (FastAPI)  |  "
            "python -m src.orchestrator.pipeline --input data/raw/ue_64_6hr.pcap  |  "
            "make test-all  (173 passed, 1 skipped, 0 failed)",
            8.5, YELLOW, wrap=True)
    footer(slide, n)


def s_c_journal(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "Journal Paper Draft — Final Results")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)

    textbox(slide, 0.3, 0.65, 9.4, 0.28,
            "Target venue: IEEE Communications Letters / Elsevier Computer Networks  |  "
            "Status: Draft complete; numbers synchronized with authoritative benchmark run (seed=42).",
            9.5, GREY, italic=True, wrap=True)

    textbox(slide, 0.3, 1.0, 9.4, 0.28, "Working Title:", 10.5, ACCENT, bold=True)
    textbox(slide, 0.3, 1.28, 9.4, 0.36,
            '"Unified Multi-Modal Anomaly Detection with Cross-Domain Correlation '
            'and LLM-Aided Explanation for 5G Networks"',
            12, GREEN, bold=True, wrap=True)

    textbox(slide, 0.3, 1.75, 9.4, 0.26, "Paper Abstract (draft — final numbers):", 10.5, ACCENT, bold=True)
    abstract = (
        "We propose a unified open-source pipeline for 5G network fault detection that "
        "jointly processes PCAP protocol traces, per-cell KPI exports, and gNB L1/L2 "
        "statistics through 18 detectors (6 per domain) covering six anomaly-detection families. "
        "A shared event correlator links anomalies across the three data domains by cell, "
        "category, and time-window, enabling cross-domain root-cause attribution that "
        "single-source methods cannot achieve. On a controlled 4-gNB, 16-cell, 64-UE synthetic "
        "testbed with a 6-hour diurnal trace and three concurrent injected fault scenarios "
        "(congestion, outage, drift), the Stats Count-Threshold method achieves "
        "Precision=1.00, Recall=1.00, F1=1.00 under the final cell-level evaluation; the "
        "Full Ensemble achieves Recall=1.00 with Precision=0.1875, F1=0.316. "
        "Holt-Winters and LSTM forecasters benchmarked (rolling-origin backtest, seed=42): "
        "HW MAE=15.54 / RMSE=21.04 / MAPE=45.67%, LSTM MAE=20.24 / RMSE=24.12 / MAPE=55.47%. "
        "Prophet implemented but not executed in the authoritative benchmark environment. "
        "End-to-end pipeline latency: 45.6 s on a CPU-only host (tracemalloc peak: ~21 MB). "
        "173 automated tests, 0 failures."
    )
    textbox(slide, 0.3, 2.03, 9.4, 1.55, abstract, 10.5, WHITE, wrap=True)

    textbox(slide, 0.3, 3.65, 9.4, 0.26, "Paper Structure:", 10.5, ACCENT, bold=True)
    sections = [
        "I.  Introduction — 5G monitoring gaps, single-source limitations, contributions listed",
        "II.  Related Work — PCAP tools, KPI threshold methods, LLM-for-telecom survey",
        "III.  System Architecture — parsers, 18-detector ensemble, EventRouter, prediction, LLM",
        "IV.  Dataset & Experimental Setup — synthetic 64-UE/16-cell/6-hr topology, injected faults",
        "V.  Results & Analysis — per-domain detection, correlation, benchmark comparison table",
        "VI.  Conclusion & Future Work — real hardware validation, online learning, LLM fine-tuning",
    ]
    for i, sec in enumerate(sections):
        ty = 3.94 + i * 0.235
        textbox(slide, 0.4, ty, 9.2, 0.22, sec, 9.5, WHITE, wrap=True)

    textbox(slide, 0.3, 5.35, 9.4, 0.22,
            "Draft location: docs/document_journalpaper/  |  Figures: all system diagrams generated",
            9, GREY, italic=True, wrap=True)
    footer(slide, n)


def s_c_references(prs, n):
    slide = new_slide(prs)
    slide_title(slide, "References")
    section_badge(slide, "Part C — Review 3", PART_C_COLOR)
    refs = [
        "3GPP TS 24.501 — Non-Access-Stratum (NAS) protocol for 5G System",
        "3GPP TS 38.413 — NG Application Protocol (NGAP)",
        "3GPP TS 38.331 — Radio Resource Control (RRC) protocol specification",
        "3GPP TS 38.473 / 38.463 / 38.423 — F1AP, E1AP, XnAP protocols",
        "3GPP TS 38.913 — Requirements for next-generation access technologies",
        "3GPP TS 38.321 — Medium Access Control (MAC) protocol specification",
        "Pedregosa et al. — Scikit-learn: Machine Learning in Python, JMLR 2011 "
        "(Isolation Forest, OC-SVM, LOF, Elliptic Envelope)",
        "Hochreiter & Schmidhuber — Long Short-Term Memory, Neural Computation 1997 (LSTM-AE)",
        "Taylor & Letham — Forecasting at Scale (Prophet), Am. Statistician 2018",
        "Holt & Winters — Exponential smoothing / damped trend (statsmodels implementation)",
        "Page — Continuous inspection schemes, Biometrika 1954 (CUSUM algorithm)",
        "Bollinger — Bollinger on Bollinger Bands, McGraw-Hill 2001",
        "Johnson, Douze & Jégou — Billion-Scale Similarity Search with GPUs (FAISS), IEEE TBDATA 2019",
        "Reimers & Gurevych — Sentence-BERT (all-MiniLM-L6-v2), EMNLP 2019",
        "Yao et al. — ReAct: Synergizing Reasoning and Acting in LMs, ICLR 2023",
        "Ollama project — Local LLM runtime (phi3:mini), github.com/ollama/ollama",
    ]
    for i, r in enumerate(refs):
        top = 0.68 + i * 0.305
        textbox(slide, 0.4, top, 9.2, 0.29, f"[{i+1}]  {r}", 9.5, WHITE, wrap=True)
    footer(slide, n)


def s_close(prs, n):
    slide = new_slide(prs)
    textbox(slide, 0.5, 1.1, 9.0, 0.7,
            "Thank You", 38, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.9, 9.0, 0.4,
            "Questions & Discussion", 17, WHITE, align=PP_ALIGN.CENTER)
    hline(slide, 2.42, ACCENT)
    textbox(slide, 0.5, 2.58, 9.0, 0.28,
            "Summary of Phase II Journey:", 13, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    summary_rows = [
        ("Review 1", PART_A_COLOR, "Novelty proposed — two Phase II research directions, 40% code"),
        ("Review 2", PART_B_COLOR, "Both novelties implemented — ReAct agent + forecast benchmark, 80% code, 83 tests"),
        ("Review 3", PART_C_COLOR, "Final validation — 64 UE / 6 hr, 173 tests passed, 100% code, journal draft prepared  |  Best F1=1.00 (Stats Count-Threshold, 3 injected fault cells)"),
    ]
    for i, (label, clr, desc) in enumerate(summary_rows):
        y = 2.96 + i * 0.5
        textbox(slide, 0.6, y, 1.5, 0.4, label, 12, clr, bold=True)
        textbox(slide, 2.2, y, 7.3, 0.4, desc, 11.5, WHITE, wrap=True)
    footer(slide, n)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)

    pg = 1
    # Cover
    s_cover(prs);  pg += 1

    # ── Part A: Phase II First Review ──────────────────────────────────────
    s_part_a_banner(prs);              pg += 1
    s_a_phase1_recap(prs, pg);        pg += 1
    s_a_novelty(prs, pg);             pg += 1
    s_a_system_ph2(prs, pg);          pg += 1
    s_a_algos(prs, pg);               pg += 1
    s_a_code_status(prs, pg);         pg += 1

    # ── Part B: Phase II Second Review ─────────────────────────────────────
    s_part_b_banner(prs);             pg += 1
    s_b_abstract(prs, pg);            pg += 1
    s_b_modified_algo(prs, pg);       pg += 1
    s_b_contribution(prs, pg);        pg += 1
    s_b_results(prs, pg);             pg += 1
    s_b_code_status(prs, pg);         pg += 1

    # ── Part C: Phase II Third Review ──────────────────────────────────────
    s_part_c_banner(prs);             pg += 1
    s_c_title_abstract(prs, pg);      pg += 1
    s_c_overall_design(prs, pg);      pg += 1
    s_c_experimental_results(prs, pg); pg += 1
    s_c_performance(prs, pg);         pg += 1
    s_c_comparison(prs, pg);          pg += 1
    s_c_code_demo(prs, pg);           pg += 1
    s_c_journal(prs, pg);             pg += 1
    s_c_references(prs, pg);          pg += 1
    s_close(prs, pg)

    out = Path("docs/Phase2_Third_Review_Comprehensive.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
