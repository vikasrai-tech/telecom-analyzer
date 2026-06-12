"""
Generate final project overview PPT — Unified Telecom Analyzer (Phase I Complete)
12 slides covering all components, 56 tests, demo flow, reviewer Q&A
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x2D, 0xC6, 0x53)
ORANGE   = RGBColor(0xFF, 0x94, 0x00)
RED      = RGBColor(0xFF, 0x4D, 0x4D)
GREY     = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE   = RGBColor(0xA0, 0x60, 0xFF)
PINK     = RGBColor(0xFF, 0x4D, 0x9F)
LIGHT_BG = RGBColor(0x14, 0x2B, 0x40)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def slide_title(slide, text):
    textbox(slide, 0.3, 0.15, 9.4, 0.5, text, 24, ACCENT, bold=True)


# ── Slide 1: Title ────────────────────────────────────────────────────

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    textbox(slide, 0.5, 1.0, 9.0, 0.9,
            "Unified Telecom Analyzer", 42, ACCENT, bold=True,
            align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.1, 9.0, 0.5,
            "End-to-End 5G Network Anomaly Detection & LLM Explanation",
            20, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.75, 9.0, 0.4,
            "M.Tech Data Science  |  PES Bangalore — Great Learning  |  Phase I Complete",
            13, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.35, 9.0, 0.4,
            "✅ 6 Protocols  ✅ 6 PCAP Detectors  ✅ 6 KPI Detectors  "
            "✅ 6 Stats Detectors  ✅ RAG + LLM  ✅ 56 Tests Passing",
            13, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.9, 9.0, 0.4,
            "Python · pyshark · scikit-learn · PyTorch · Prophet · FAISS · Ollama · Streamlit · FastAPI",
            12, GREY, align=PP_ALIGN.CENTER)


# ── Slide 2: Architecture ─────────────────────────────────────────────

def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "🏗️  System Architecture — 5 Layers")

    layers = [
        ("INPUT",
         "PCAP (.pcap/.pcapng)  ·  DU/CU Stats (srsRAN / OAI / NIST CSV/Parquet)  ·  KPI Time-series (Excel/CSV)",
         ACCENT, 0.75),
        ("PARSING  —  TS 24.501 / 38.413 / 38.331 / 38.473 / 38.463 / 38.423",
         "NAS  ·  NGAP  ·  RRC  ·  F1AP  ·  E1AP  ·  XnAP  ·  KPI Parser  ·  Stats Parser (srsRAN/OAI/NIST)",
         GREEN, 1.55),
        ("DETECTION  —  18 detectors total",
         "PCAP: IF · Statistical · OC-SVM · LOF · Elliptic Env. · LSTM-AE    "
         "KPI: Threshold · Peer Comp. · Trend · IQR · CUSUM · Bollinger    "
         "Stats: same 6 methods on L1/L2 counters",
         ORANGE, 2.45),
        ("RAG + LLM  —  FAISS · MiniLM · Ollama phi3:mini  +  Feedback Loop + Prediction",
         "Query build → FAISS top-5 retrieval → Ollama prompt → structured JSON output  "
         "|  Feedback store → nightly retrainer  |  LSTM + Prophet 4h-ahead prediction",
         PURPLE, 3.35),
        ("OUTPUT",
         "Streamlit Dashboard (http://localhost:8501)  ·  FastAPI REST API (:8000)  ·  CLI Pipeline Runner",
         PINK, 4.2),
    ]

    for label, items, color, top in layers:
        textbox(slide, 0.3, top,      2.5, 0.3, label, 9,  color, bold=True)
        textbox(slide, 0.3, top+0.28, 9.2, 0.45, items, 10, WHITE, wrap=True)


# ── Slide 3: Protocols ────────────────────────────────────────────────

def add_protocol_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "📡  3GPP Protocols Parsed — Full 5G-NR Stack")

    protocols = [
        ("NAS — TS 24.501",
         "Registration · Auth · Security Mode · PDU Session Establishment · Deregistration",
         "UE ↔ AMF  (N1)"),
        ("NGAP — TS 38.413",
         "InitialContextSetup · PDUSession Setup/Release · UEContextRelease · Handover · Paging",
         "gNB-CU ↔ AMF  (N2)"),
        ("RRC — TS 38.331",
         "Setup · Reestablishment · Reconfiguration · Release · UE Capability · Measurement",
         "UE ↔ gNB  (Uu)"),
        ("F1AP — TS 38.473",
         "F1 Setup · UE Context Setup/Release · DL/UL RRC Transfer · Initial UL RRC",
         "gNB-DU ↔ gNB-CU-CP  (F1)"),
        ("E1AP — TS 38.463",
         "CU-UP Setup · Bearer Context Setup/Modify/Release · Data Notification",
         "gNB-CU-CP ↔ gNB-CU-UP  (E1)"),
        ("XnAP — TS 38.423",
         "Xn Setup · Handover Request · UE Context Release · SN Addition (MR-DC)",
         "gNB ↔ gNB  (Xn)"),
    ]

    for i, (proto, procs, iface) in enumerate(protocols):
        top = 0.85 + i * 0.73
        textbox(slide, 0.3, top, 2.5, 0.3, proto, 12, GREEN, bold=True)
        textbox(slide, 2.9, top, 5.0, 0.3, procs, 10, WHITE)
        textbox(slide, 8.0, top, 1.8, 0.3, iface,  9, GREY,
                align=PP_ALIGN.RIGHT)


# ── Slide 4: PCAP Detection ───────────────────────────────────────────

def add_pcap_detection_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "⚠️  PCAP Anomaly Detection — 6-Detector Ensemble")

    detectors = [
        ("1", "Isolation Forest",
         "Unsupervised / tree",
         "Anomalies isolate in fewer tree splits. No distribution assumption. O(n log n). Best for global outliers."),
        ("2", "Statistical (Threshold+Cascade)",
         "Rule-based / domain",
         "Encodes 3GPP SLA directly. 100% interpretable. Catches congestion chains, timeout storms, NAS→NGAP cascades."),
        ("3", "One-Class SVM",
         "Kernel / boundary",
         "Non-linear normal region in RKHS. Catches anomalies near but outside the normal boundary."),
        ("4", "LOF",
         "Density / local",
         "Local Outlier Factor — catches cells 3σ below their local peer cluster even when globally within range."),
        ("5", "Elliptic Envelope",
         "Mahalanobis / Gaussian",
         "Fastest interpretable baseline. Fits multivariate Gaussian; sanity-checks all ML methods."),
        ("6", "LSTM Autoencoder",
         "Deep learning / sequence",
         "ONLY method that catches procedure-ORDER violations (e.g., PDU Session before Authentication)."),
    ]

    for i, (num, name, dtype, why) in enumerate(detectors):
        top = 0.82 + i * 0.63
        textbox(slide, 0.3,  top, 0.25, 0.42, num,   14, ACCENT, bold=True)
        textbox(slide, 0.6,  top, 2.2,  0.42, name,  11, WHITE,  bold=True)
        textbox(slide, 2.85, top, 1.6,  0.42, dtype, 10, ORANGE)
        textbox(slide, 4.5,  top, 5.2,  0.42, why,   10, GREY,   wrap=True)


# ── Slide 5: KPI Detection ────────────────────────────────────────────

def add_kpi_detection_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "📊  KPI + L1/L2 Stats Detection — 6-Method Ensemble")

    methods = [
        ("1", "Threshold Violation",
         "Rule-based / point",
         "Per-row check vs warning/critical thresholds. Encodes 3GPP/ITU-T SLA directly. Zero false positives for known limits."),
        ("2", "Peer Comparison",
         "Statistical / cross-cell",
         "Z-score vs fleet mean. Catches underperforming cells still within threshold but 3σ below peers."),
        ("3", "Trend (Linear Regression)",
         "Temporal / slope",
         "Slope > 0.5 unit/hr = actionable degradation even if current value looks 'green'. Early-warning signal."),
        ("4", "IQR (Tukey Fence)",
         "Robust / distribution-free",
         "No Gaussian assumption. Works on skewed KPIs — PRB right-skewed, CQI left-skewed during interference."),
        ("5", "CUSUM",
         "Sequential / change-point",
         "Accumulates small deviations. Catches RRC SR 99.5%→98.2% drift 5+ hours before threshold fires."),
        ("6", "Bollinger Bands",
         "Rolling envelope / burst",
         "Catches single-timestamp spikes invisible to trend methods. Rolling window gives local context."),
    ]

    textbox(slide, 0.3, 0.7, 9.2, 0.25,
            "Same 6 methods applied to both KPI time-series (Excel/CSV) and DU/CU L1/L2 statistics (srsRAN · OAI · NIST)",
            10, GREY)

    for i, (num, name, dtype, why) in enumerate(methods):
        top = 1.05 + i * 0.62
        textbox(slide, 0.3,  top, 0.25, 0.42, num,   14, ACCENT, bold=True)
        textbox(slide, 0.6,  top, 2.2,  0.42, name,  11, WHITE,  bold=True)
        textbox(slide, 2.85, top, 1.8,  0.42, dtype, 10, ORANGE)
        textbox(slide, 4.7,  top, 5.0,  0.42, why,   10, GREY,   wrap=True)


# ── Slide 6: LLM / RAG ───────────────────────────────────────────────

def add_llm_rag_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "🤖  LLM Explainer — RAG over 3GPP Specifications")

    flow = [
        ("1  Anomaly input",
         "procedure + severity + failure_causes + evidence dict",
         "from PCAP / KPI / Stats detector"),
        ("2  Query build",
         "natural-language query from anomaly fields",
         "_build_query()"),
        ("3  Retrieve",
         "FAISS cosine search → top-5 3GPP spec chunks",
         "all-MiniLM-L6-v2 (80 MB), 38 curated chunks"),
        ("4  Prompt",
         "anomaly + retrieved spec context → structured JSON prompt",
         "_build_prompt()  |  temperature=0.2"),
        ("5  LLM call",
         "Ollama: phi3:mini → llama3.2:1b → mistral:7b  (preference chain)",
         "Runs fully LOCAL — no cloud API, no data leaves premises"),
        ("6  Rule-based fallback",
         "template explanation built from retrieved chunks",
         "fires automatically if Ollama unavailable"),
        ("7  Output",
         "hypothesis · severity · citations (spec §) · investigation_hints",
         "dashboard + REST API response"),
    ]

    for i, (step, desc, note) in enumerate(flow):
        top = 0.82 + i * 0.58
        textbox(slide, 0.3, top, 1.4,  0.42, step, 12, ACCENT, bold=True)
        textbox(slide, 1.8, top, 4.8,  0.42, desc, 11, WHITE)
        textbox(slide, 6.7, top, 3.0,  0.42, note,  9, GREY,   wrap=True)

    textbox(slide, 0.3, 5.0, 9.4, 0.28,
            "Knowledge base: 38 curated chunks — TS 24.501 · 38.413 · 38.331 · 38.473 · 38.463 · 38.423 · 38.913",
            10, ORANGE)


# ── Slide 7: MLOps / Feedback ─────────────────────────────────────────

def add_mlops_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "🔄  MLOps Pipeline — Feedback Loop + Prediction Layer")

    # Feedback loop
    textbox(slide, 0.3, 0.72, 4.5, 0.3, "Feedback Loop (Tier 5)", 14, GREEN, bold=True)
    feedback_items = [
        "Engineer marks each anomaly:  👍 Correct / 👎 False Positive / ❓ Uncertain",
        "Feedback stored to JSONL log  (src/feedback/store.py)",
        "Nightly retrainer reads FP rate per detector",
        "Adjusts contamination (IF/SVM/LOF/EE), sensitivity (Statistical), threshold (LSTM)",
        "Formula:  new_param = base + α × (fp_rate − target_fp_rate)  [clamped 0.02–0.40]",
        "Cron: 0 2 * * *  cd /mnt/e/telecom-analyzer && make retrain",
    ]
    for j, item in enumerate(feedback_items):
        textbox(slide, 0.5, 1.05 + j * 0.32, 4.3, 0.3, f"• {item}", 10, WHITE, wrap=True)

    # Prediction layer
    textbox(slide, 5.0, 0.72, 4.7, 0.3, "Prediction Layer — 4h Ahead", 14, PURPLE, bold=True)
    pred_items = [
        "Prophet — Bayesian structural time-series",
        "  Handles seasonality · trends · missing data",
        "  Best for KPIs with daily/weekly cycles (PRB, throughput)",
        "LSTM Sequence Model — nonlinear dependencies",
        "  Catches abrupt non-seasonal shifts (HARQ NACK, BLER)",
        "Output: anomaly dicts tagged  state='predicted'  + lead_time_h",
        "Dashboard: 🔮 Prediction Layer section with horizon slider",
    ]
    for j, item in enumerate(pred_items):
        textbox(slide, 5.2, 1.05 + j * 0.32, 4.5, 0.3, f"• {item}", 10, WHITE, wrap=True)

    textbox(slide, 0.3, 5.05, 9.4, 0.28,
            "make retrain  |  scripts/nightly_retrain.py  |  data/models/retrain_config.json stores tuned params",
            10, GREY)


# ── Slide 8: Output Interfaces ───────────────────────────────────────

def add_api_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "🌐  Output Interfaces — Dashboard · REST API · CLI")

    interfaces = [
        ("Streamlit Dashboard", "make dev  →  http://localhost:8501",
         ["Upload PCAP, Stats CSV, or KPI Excel in one interface",
          "6 protocol tabs · failure cause breakdown · IE inspector",
          "6-detector method comparison matrix + rationale expanders",
          "LLM explanation with 3GPP spec citations",
          "Prediction panel (LSTM + Prophet) · Feedback buttons · Event Router"]),
        ("FastAPI REST API", "make api  →  http://localhost:8000/docs",
         ["POST /analyze/pcap  — upload PCAP, returns full anomaly JSON",
          "POST /analyze/kpi   — upload KPI file, returns anomaly JSON",
          "GET  /health        — Ollama status · model · version"]),
        ("CLI Pipeline Runner",
         "python -m src.orchestrator.pipeline --input file.pcap [--no-llm] [--output report.json]",
         ["Auto-detects PCAP / KPI / Stats from file extension",
          "--no-llm skips LLM (fast mode for CI)",
          "--output writes full JSON report for batch processing"]),
    ]

    top = 0.75
    for name, cmd, bullets in interfaces:
        textbox(slide, 0.3, top,        9.2, 0.3, f"▶  {name}  —  {cmd}", 12, GREEN, bold=True)
        for j, b in enumerate(bullets):
            textbox(slide, 0.7, top + 0.32 + j * 0.28, 9.0, 0.28, f"• {b}", 10, WHITE)
        top += 0.32 + len(bullets) * 0.28 + 0.22


# ── Slide 9: Demo Flow ───────────────────────────────────────────────

def add_demo_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "🎬  Demo Flow — Live Walkthrough (15 min)")

    steps = [
        ("Step 1  (2 min)", "Architecture intro",
         "Show README diagram. Explain 3 data paths, 18 detectors, LLM explainer.",
         ACCENT),
        ("Step 2  (5 min)", "PCAP Demo  →  upload test_5g_full.pcap",
         "Parsed Summary → Procedure Counters → Failure Cause Breakdown → "
         "6-Detector Ensemble → Method Comparison Matrix → LLM Explanation",
         GREEN),
        ("Step 3  (3 min)", "KPI Demo  →  upload 5G_Network_KPI_Sample.xlsx",
         "KPI Health Summary heatmap → Trend Explorer → Per-Cell Breakdown → "
         "6-Method Ensemble → Feedback buttons (👍/👎)",
         ORANGE),
        ("Step 4  (2 min)", "Stats Demo  →  upload srsran_stats.csv",
         "L1/L2 metrics (BLER, HARQ, SNR, PRB) → anomaly detection at physical layer",
         PURPLE),
        ("Step 5  (2 min)", "MLOps + REST API",
         "Show Prediction Layer (4h forecast) · Feedback History · "
         "Open http://localhost:8000/docs → POST /analyze/pcap live",
         PINK),
        ("Step 6  (1 min)", "Test suite",
         "Run  make test  in terminal → show 56 passing tests",
         GREY),
    ]

    for i, (step, title, detail, color) in enumerate(steps):
        top = 0.78 + i * 0.73
        textbox(slide, 0.3,  top,       1.2, 0.28, step,   10, color, bold=True)
        textbox(slide, 1.55, top,       3.5, 0.28, title,  11, WHITE, bold=True)
        textbox(slide, 1.55, top+0.3,   8.1, 0.32, detail,  9, GREY,  wrap=True)


# ── Slide 10: Test Results ────────────────────────────────────────────

def add_test_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "✅  Test Results — 56 Tests, All Passing")

    suites = [
        ("test_walking_skeleton.py  (4 tests)",
         ["parse_pcap_stub keys", "detect_anomalies_stub list",
          "explain_anomaly_stub citations", "stub pipeline end-to-end"],
         GREEN, 0),
        ("test_real_pipeline.py  (16 tests)",
         ["Real PCAP parser — keys, procedures, schema",
          "6 PCAP detectors — correct dict schema",
          "KPI parser + 6 KPI detectors validated",
          "RAG retriever — chunks + relevance ordering",
          "LLM explainer — rule-based output schema",
          "run_pcap_pipeline() and run_kpi_pipeline() end-to-end",
          "run_pipeline() auto-detects file type"],
         GREEN, 1),
        ("test_event_router.py  (8 tests)",
         ["ingest + summary + correlated events",
          "cross-source correlation, top cells, event log"],
         GREEN, 2),
        ("test_feedback_prediction.py  (12 tests)",
         ["feedback save/load/stats", "LSTM + Prophet prediction schema"],
         GREEN, 3),
        ("test_retrainer.py  (8 tests)",
         ["FP rate calculation", "param adjustment formula", "dry-run mode"],
         GREEN, 4),
        ("test_walking_skeleton.py  (8 tests)",
         ["stats parser — srsRAN/OAI/NIST", "6 stats detectors schema"],
         GREEN, 5),
    ]

    col_w   = 4.5
    columns = [(0.3, 0), (5.0, 3)]

    for title, tests, color, idx in suites:
        col_idx = 0 if idx < 3 else 1
        base_x, _ = columns[col_idx]
        row = idx if idx < 3 else idx - 3
        top = 0.78 + row * 1.45

        textbox(slide, base_x, top, col_w, 0.3, title, 11, ACCENT, bold=True)
        for j, t in enumerate(tests):
            textbox(slide, base_x+0.2, top+0.32+j*0.28, col_w-0.2, 0.28,
                    f"✅  {t}", 10, color)


# ── Slide 11: Reviewer Q&A ───────────────────────────────────────────

def add_reviewer_qa_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title(slide, "🎯  Reviewer Q&A — Anticipated Questions")

    qa = [
        ("Why 6 detectors instead of 1?",
         "No single method catches all anomaly types. IF handles global; LOF handles local density. "
         "Statistical is deterministic; LSTM catches sequence violations. "
         "Multi-detector agreement = higher confidence, lower false positive rate."),
        ("Why Ollama (local LLM) not OpenAI/cloud?",
         "Telecom data is sensitive — operator configs, failure rates, cell IDs cannot leave premises. "
         "Ollama runs phi3:mini fully local. Rule-based fallback ensures system works even without GPU."),
        ("How does the RAG know 3GPP specs?",
         "38 curated spec excerpts from TS 24.501/38.413/38.331/38.473/38.463/38.423 are embedded "
         "with all-MiniLM-L6-v2 into a FAISS index. Each anomaly retrieves top-5 relevant chunks by cosine similarity."),
        ("What is the feedback loop?",
         "Engineers mark TP/FP/Uncertain on each anomaly card. Nightly cron reads FP rate per detector "
         "and adjusts contamination/sensitivity. Formula: new_param = base + α×(fp_rate − target). "
         "This prevents false positive accumulation over time."),
        ("What is the detection latency?",
         "PCAP: ~5s for 500 procedures · KPI: ~3s for 400 rows · Stats: ~1s for 300 rows. "
         "Rule-based + IF are fastest; LSTM adds ~1s. --no-llm flag gives sub-2s total."),
    ]

    for i, (q, a) in enumerate(qa):
        top = 0.78 + i * 0.88
        textbox(slide, 0.3, top,      9.2, 0.3,  f"Q:  {q}", 12, ORANGE, bold=True)
        textbox(slide, 0.5, top+0.3,  9.0, 0.48, f"A:  {a}",  10, WHITE, wrap=True)


# ── Slide 12: Summary ─────────────────────────────────────────────────

def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    textbox(slide, 0.5, 0.7, 9.0, 0.7,
            "Phase I — Complete ✅", 38, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.5, 9.0, 0.3,
            "M.Tech Data Science  |  PES Bangalore — Great Learning",
            14, GREY, align=PP_ALIGN.CENTER)

    items = [
        ("Parsers",     "6 × 3GPP protocols (NAS · NGAP · RRC · F1AP · E1AP · XnAP) + KPI + Stats (srsRAN/OAI/NIST)"),
        ("Detection",   "18 detectors — 6 PCAP + 6 KPI + 6 Stats  (Isolation Forest through Bollinger Bands)"),
        ("LLM",         "RAG pipeline: 38 spec chunks · FAISS · MiniLM · Ollama phi3:mini + rule-based fallback"),
        ("MLOps",       "Feedback loop → nightly retrainer + LSTM + Prophet 4h-ahead prediction layer"),
        ("Interfaces",  "Streamlit dashboard · FastAPI REST API · CLI orchestrator · Event Router"),
        ("Tests",       "56 tests passing — walking skeleton + real integration + event + feedback + retrainer"),
        ("Docs",        "4 PPTs (PCAP methods · KPI methods · Stats · this deck) + Study Guide + SETUP"),
    ]

    for i, (label, detail) in enumerate(items):
        top = 1.95 + i * 0.44
        textbox(slide, 1.0,  top, 1.5,  0.38, f"✅  {label}", 13, ACCENT, bold=True)
        textbox(slide, 2.6,  top, 7.0,  0.38, detail,         12, WHITE)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)

    add_title_slide(prs)
    add_architecture_slide(prs)
    add_protocol_slide(prs)
    add_pcap_detection_slide(prs)
    add_kpi_detection_slide(prs)
    add_llm_rag_slide(prs)
    add_mlops_slide(prs)
    add_api_slide(prs)
    add_demo_flow_slide(prs)
    add_test_results_slide(prs)
    add_reviewer_qa_slide(prs)
    add_summary_slide(prs)

    out = Path("docs/Project_Overview_Final.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
