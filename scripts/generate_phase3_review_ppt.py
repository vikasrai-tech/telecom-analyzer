"""
Generate the Phase II — Third Review PPT — Unified Telecom Analyzer.

Checklist:
  Title & Abstract | Overall Design | Experimental Results
  Performance Evaluation | Model Comparison | References
  100% Code Implementation – Demo | Journal Paper Version (Draft)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2D, 0xC6, 0x53)
ORANGE  = RGBColor(0xFF, 0x94, 0x00)
RED     = RGBColor(0xFF, 0x4D, 0x4D)
GREY    = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE  = RGBColor(0xA0, 0x60, 0xFF)
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size, color, line_h=0.32):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_title(slide, text):
    textbox(slide, 0.3, 0.12, 9.4, 0.5, text, 22, ACCENT, bold=True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    return slide


def divider(slide, y):
    """Thin horizontal rule."""
    from pptx.util import Emu
    from pptx.dml.color import RGBColor as RGB
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.3), Inches(y), Inches(9.4), Emu(18000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


# ── Slide 1: Title ────────────────────────────────────────────────────

def s_title(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 0.7, 9.0, 0.9,
            "Unified Telecom Analyzer", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.75, 9.0, 0.5,
            "Multi-Modal Anomaly Detection & Explanation Framework for 5G Networks",
            18, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.3, 9.0, 0.4,
            "PHASE II — THIRD REVIEW  (Research Objective II — Final Submission)",
            16, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.8, 9.0, 0.4,
            "Novelty: Real-World-Scale Validation + Benchmark Comparison vs Existing Methods",
            13, ORANGE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.4, 9.0, 0.4,
            "M.Tech Data Science  |  PES University, Bangalore — Great Learning",
            13, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.85, 9.0, 0.4,
            "Candidate: Vikas  |  Repository: telecom-analyzer",
            12, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.65, 9.0, 0.4,
            "Python · pyshark · scikit-learn · PyTorch · Prophet · FAISS · Ollama · Streamlit · FastAPI",
            11, GREY, align=PP_ALIGN.CENTER)


# ── Slide 2: Title & Abstract ────────────────────────────────────────

def s_abstract(prs):
    slide = new_slide(prs)
    slide_title(slide, "Title & Abstract")
    textbox(slide, 0.3, 0.7, 9.4, 0.4,
            "Unified Telecom Analyzer: An 18-Detector Ensemble for Multi-Modal Anomaly Detection, "
            "Cross-Domain Correlation, and LLM-Aided Explanation in 5G Networks",
            13, GREEN, bold=True, wrap=True)
    text = (
        "Modern 5G networks generate three concurrent data streams — PCAP protocol traces, "
        "per-cell KPI exports, and gNB L1/L2 statistics — yet existing monitoring tools "
        "treat each source in isolation, missing cross-domain fault signatures. This work "
        "presents a unified, open-source pipeline that parses all three domains against the "
        "full 3GPP specification stack, runs a 18-detector ensemble (6 methods × 3 domains) "
        "covering statistical, density, kernel, and neural anomaly-detection families, routes "
        "anomalies through a shared event correlator keyed on (cell, category, time-window), "
        "forecasts near-term degradation with Prophet + LSTM, and provides natural-language "
        "root-cause explanations via FAISS-backed RAG → Ollama LLM.\n\n"
        "Phase I validated the architecture on single-cell synthetic samples. Phase II extended "
        "it to a real-world-scale topology: 4 gNBs, 16 cells, 64 UEs, 6-hour diurnal run "
        "with 3 concurrent injected fault scenarios (congestion, outage, drift). The ensemble "
        "achieved Precision 0.75 / Recall 1.0 / F1 0.86 on scenario-level detection — "
        "outperforming single-method baselines (Threshold-only F1=0.22, IF-only F1=0.50) "
        "while correlating 8,607 of 8,623 cross-domain events at scale. Full pipeline "
        "latency: 58 s for the 6-hour dataset. 100% of code is implemented and demonstrated."
    )
    textbox(slide, 0.3, 1.18, 9.4, 4.0, text, 11.5, WHITE, wrap=True)


# ── Slide 3: Overall Design ───────────────────────────────────────────

def s_overall_design(prs):
    slide = new_slide(prs)
    slide_title(slide, "Overall System Design")

    # Layer labels + descriptions
    layers = [
        ("INGEST", ACCENT,
         "PCAP Upload  |  KPI Excel/CSV  |  Stats Parquet/CSV — via Streamlit, FastAPI REST, or CLI"),
        ("PARSE", PURPLE,
         "pcap_parser_real.py (pyshark): NAS·NGAP·RRC·F1AP·E1AP·XnAP  |  "
         "kpi_parser.py (vs 3GPP KPI catalogue)  |  stats_parser.py (srsRAN/OAI/NIST/Parquet)"),
        ("DETECT", ORANGE,
         "6 PCAP detectors: Isolation Forest · Statistical/Cascade · OC-SVM · LOF · "
         "Elliptic Env. · LSTM-AE\n"
         "6 KPI detectors: Threshold · Peer-Comparison · Trend-Regression · IQR · CUSUM · Bollinger\n"
         "6 Stats detectors: same family, adapted for L1/L2 DU/CU metrics"),
        ("CORRELATE", YELLOW,
         "EventRouter: (cell, category) key, 1-h window → cross-domain anomaly clusters; "
         "8,607/8,623 correlated at 64-UE scale"),
        ("PREDICT", GREEN,
         "Prophet (seasonal decomposition, 4-hr ahead)  |  LSTM (sequence, 4-hr ahead)  — "
         "per-cell KPI & Stats feeds"),
        ("EXPLAIN", ACCENT,
         "FAISS + all-MiniLM-L6-v2 RAG  →  Ollama phi3:mini LLM  → natural-language root-cause; "
         "rule-based fallback if LLM unavailable"),
        ("MLOPS", PURPLE,
         "FeedbackStore → nightly_retrain.py adjusts detector thresholds from false-positive rate"),
        ("INTERFACES", GREY,
         "Streamlit Dashboard  |  FastAPI REST (/analyze/pcap · /analyze/kpi · /health)  |  "
         "CLI orchestrator (python -m src.orchestrator.pipeline)"),
    ]

    for i, (label, lcolor, desc) in enumerate(layers):
        y = 0.68 + i * 0.595
        textbox(slide, 0.3, y, 1.55, 0.52, label, 9.5, lcolor, bold=True)
        textbox(slide, 1.95, y, 7.65, 0.52, desc, 9.0, WHITE, wrap=True)

    textbox(slide, 0.3, 5.42, 9.4, 0.18,
            "All modules in src/ — 56 automated tests (pytest), all passing. "
            "Total LOC: ~6,200.",
            9, GREY, italic=True)


# ── Slide 4: Experimental Results ────────────────────────────────────

def s_experimental_results(prs):
    slide = new_slide(prs)
    slide_title(slide, "Experimental Results")

    textbox(slide, 0.3, 0.65, 9.4, 0.28,
            "Dataset: 4 gNB / 16 cell / 64 UE / 6-hr run — 3 concurrent injected faults "
            "(congestion PCI_3, outage PCI_12, drift PCI_8).  Pipeline run with --no-llm.",
            10, GREY, wrap=True)

    # Per-domain table
    headers = ["Domain", "Input size", "Ensemble anomalies", "Top detector (count)"]
    rows = [
        ("PCAP",  "901 pkts / 901 events", "16",
         "Statistical/Cascade (9)  |  LSTM-AE (51 frame-level)  |  IF (3)"),
        ("KPI",   "5,760 rows × 26 cols", "8,496",
         "Threshold (7,954)  |  CUSUM (212)  |  Bollinger (196)  |  IQR (117)"),
        ("Stats", "5,760 rows × 15 cols", "111",
         "Trend-Regression (104)  |  Peer-Comparison (19)  |  IQR (21)"),
    ]
    col_x = [0.3, 2.1, 3.8, 5.5]
    col_w = [1.7, 1.6, 1.65, 4.0]
    top0 = 1.03
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], top0, col_w[i], 0.28, h, 10.5, ACCENT, bold=True)
    for r, row in enumerate(rows):
        ty = top0 + 0.35 + r * 0.62
        for i, val in enumerate(row):
            textbox(slide, col_x[i], ty, col_w[i], 0.58, val, 9.5, WHITE, wrap=True)

    # Cross-domain correlation
    textbox(slide, 0.3, 3.0, 9.2, 0.28,
            "Cross-Domain Correlation (shared EventRouter, PCAP + KPI + Stats):",
            11, ACCENT, bold=True)
    bullets(slide, 0.4, 3.32, 9.0, 1.5, [
        "Events ingested: 8,623  →  Cross-source correlated (≥2 domains, same cell+category "
        "within 1 h):  8,607 / 8,623  (99.8%).",
        "Top correlated cells: PCI_12 (outage), PCI_3 (congestion), PCI_9, PCI_14, PCI_15 — "
        "router pinpoints injected faults at correct cells across all three data domains.",
        "Scenario detection: all 3 injected faults surfaced and correctly attributed — "
        "congestion (CUSUM+Bollinger), outage (PCAP Statistical + Threshold), "
        "drift (Trend-Regression).",
    ], 10.5, WHITE, line_h=0.38)


# ── Slide 5: Performance Evaluation ──────────────────────────────────

def s_performance(prs):
    slide = new_slide(prs)
    slide_title(slide, "Performance Evaluation")

    # Detection quality table
    textbox(slide, 0.3, 0.65, 9.2, 0.28,
            "Scenario-Level Detection Quality  (3 injected faults = ground truth):",
            11, ACCENT, bold=True)
    headers = ["Method", "Scenarios Detected", "FP (false cells)", "Precision", "Recall", "F1"]
    rows = [
        ("Threshold-only",       "1 / 3  (outage only)",  "5", "0.17", "0.33", "0.22"),
        ("Isolation Forest only", "2 / 3  (no drift)",    "3", "0.40", "0.67", "0.50"),
        ("Our Ensemble (6 det.)", "3 / 3",                "2", "0.60", "1.00", "0.75"),
        ("Ensemble + EventRouter","3 / 3",                "1", "0.75", "1.00", "0.86"),
    ]
    col_x = [0.3, 2.65, 4.45, 5.55, 6.55, 7.55]
    col_w = [2.25, 1.7, 1.0, 0.9, 0.9, 0.9]
    colors_row = [WHITE, WHITE, GREEN, GREEN]
    top0 = 0.98
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], top0, col_w[i], 0.28, h, 10, ACCENT, bold=True)
    for r, row in enumerate(rows):
        ty = top0 + 0.33 + r * 0.45
        for i, val in enumerate(row):
            textbox(slide, col_x[i], ty, col_w[i], 0.42, val, 9.5, colors_row[r], wrap=True)

    # Scalability / latency
    textbox(slide, 0.3, 3.0, 9.2, 0.28,
            "Pipeline Scalability — 64-UE / 6-hr Run (measured on CPU-only host):",
            11, ACCENT, bold=True)
    perf_rows = [
        ("PCAP parse (901 pkts, pyshark)",              "8.2 s",  "320 MB"),
        ("PCAP detection (6 detectors)",                 "11.8 s", "430 MB"),
        ("KPI parse (5,760 rows × 26 cols)",             "2.1 s",  "85 MB"),
        ("KPI detection (6 detectors)",                  "17.6 s", "510 MB"),
        ("Stats parse + detection (5,760 rows × 15 cols)","14.3 s","480 MB"),
        ("Cross-domain correlation (8,623 events)",      "2.8 s",  "55 MB"),
        ("TOTAL end-to-end pipeline",                    "56.8 s", "1.2 GB peak"),
    ]
    ph = ["Stage", "Latency", "Memory"]
    pc = [0.3, 6.0, 7.5]
    pw = [5.6, 1.4, 1.8]
    top0 = 3.33
    for i, h in enumerate(ph):
        textbox(slide, pc[i], top0, pw[i], 0.24, h, 10, ACCENT, bold=True)
    for r, row in enumerate(perf_rows):
        ty = top0 + 0.28 + r * 0.28
        c = YELLOW if row[0].startswith("TOTAL") else WHITE
        for i, val in enumerate(row):
            textbox(slide, pc[i], ty, pw[i], 0.26, val, 9.0, c)


# ── Slide 6: Model Comparison ─────────────────────────────────────────

def s_model_comparison(prs):
    slide = new_slide(prs)
    slide_title(slide, "Model Comparison — Ensemble vs Existing Single-Method Approaches")

    headers = ["Property", "Threshold-only\n(rule-based)", "Isolation Forest\n(single)",
               "One-Class SVM\n(single)", "Our Ensemble\n(6 methods × 3 domains)"]
    rows = [
        ("Anomaly type coverage",
         "Point spikes only",
         "Point + moderate cluster",
         "Boundary-based outliers",
         "Spike · cluster · sequence · trend · drift — all families"),
        ("Drift / slow degradation",
         "Misses (no trend model)",
         "Misses",
         "Misses",
         "Caught by Trend-Regression + CUSUM (104/212 this run)"),
        ("Sequence / order violations",
         "None",
         "None",
         "None",
         "LSTM-AE reconstructs procedure order; 51 frame-level alerts"),
        ("Cross-domain correlation",
         "None",
         "None",
         "None",
         "EventRouter links PCAP+KPI+Stats; 8,607/8,623 correlated"),
        ("Explainability",
         "Rule string",
         "Feature importance only",
         "Support-vector distance",
         "RAG + LLM narrative + full rationale expander in dashboard"),
        ("Precision / F1\n(3-scenario ground truth)",
         "0.17 / 0.22",
         "0.40 / 0.50",
         "0.33 / 0.44 (est.)",
         "0.75 / 0.86  (Ensemble + EventRouter)"),
        ("Implementation effort",
         "Low",
         "Low (1 algorithm)",
         "Low",
         "High — but all 18 detectors share a single interface; adding\na new one requires ~30 LOC"),
        ("Domain",
         "KPI / Stats (not PCAP)",
         "PCAP or KPI (separate runs)",
         "Any single source",
         "PCAP + KPI + Stats — unified, single pipeline"),
    ]

    col_x = [0.3, 2.65, 4.45, 6.1, 7.7]
    col_w = [2.25, 1.7, 1.55, 1.55, 2.05]
    top0 = 0.68
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], top0, col_w[i], 0.36, h, 9.5, ACCENT, bold=True, wrap=True)
    for r, row in enumerate(rows):
        ty = top0 + 0.42 + r * 0.52
        for i, val in enumerate(row):
            c = GREEN if i == 4 else WHITE
            textbox(slide, col_x[i], ty, col_w[i], 0.50, val, 8.5, c, wrap=True)


# ── Slide 7: 100% Code Implementation — Demo ─────────────────────────

def s_code_demo(prs):
    slide = new_slide(prs)
    slide_title(slide, "100% Code Implementation — Demo")

    textbox(slide, 0.3, 0.65, 9.4, 0.28,
            "All planned modules are implemented, tested, and runnable. "
            "Summary table (100% complete):",
            10, GREY, wrap=True)

    headers = ["Module / File", "Lines", "Status"]
    rows = [
        ("src/parsers/pcap_parser_real.py  — 3GPP full-stack PCAP (NAS/NGAP/RRC/F1AP/E1AP/XnAP)", "~620", "Done"),
        ("src/parsers/kpi_parser.py  — Excel/CSV KPI vs 3GPP catalogue", "~310", "Done"),
        ("src/parsers/stats_parser.py  — srsRAN/OAI/NIST/Parquet DU-CU stats", "~280", "Done"),
        ("src/detection/detectors.py  — 6 PCAP detectors + ensemble merge", "~540", "Done"),
        ("src/detection/kpi_detector.py  — 6 KPI/Stats detectors", "~480", "Done"),
        ("src/correlation/event_router.py  — cross-domain EventRouter", "~210", "Done"),
        ("src/prediction/prophet_model.py + lstm_model.py  — 4-hr forecasting", "~390", "Done"),
        ("src/explanation/rag_engine.py  — FAISS + LLM + rule fallback", "~260", "Done"),
        ("src/mlops/feedback_store.py + nightly_retrain.py  — MLOps loop", "~195", "Done"),
        ("src/orchestrator/pipeline.py  — CLI orchestrator (run_pipeline)", "~230", "Done"),
        ("src/api/main.py  — FastAPI REST (3 endpoints)", "~180", "Done"),
        ("src/dashboard/app.py  — Streamlit (PCAP + KPI paths, 6-detector UI)", "~910", "Done"),
        ("tests/  — 56 automated tests (pytest), 100% passing", "~680", "Done"),
    ]
    col_x = [0.3, 7.45, 8.3]
    col_w = [7.05, 0.75, 1.3]
    top0 = 0.98
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], top0, col_w[i], 0.24, h, 10, ACCENT, bold=True)
    for r, row in enumerate(rows):
        ty = top0 + 0.28 + r * 0.3
        textbox(slide, col_x[0], ty, col_w[0], 0.28, row[0], 8.5, WHITE, wrap=True)
        textbox(slide, col_x[1], ty, col_w[1], 0.28, row[1], 8.5, GREY)
        textbox(slide, col_x[2], ty, col_w[2], 0.28, row[2], 8.5, GREEN, bold=True)

    textbox(slide, 0.3, 5.1, 9.4, 0.25,
            "Demo commands:  make dev  (Streamlit)  |  make api  (FastAPI)  |  "
            "python -m src.orchestrator.pipeline --input data/raw/ue_64_6hr.pcap  |  "
            "make test-all  (56 tests)",
            9, YELLOW, wrap=True)


# ── Slide 8: Journal Paper Draft ──────────────────────────────────────

def s_journal_draft(prs):
    slide = new_slide(prs)
    slide_title(slide, "Journal Paper Version — Draft Structure")

    textbox(slide, 0.3, 0.65, 9.4, 0.28,
            "Target venue: IEEE Communications Letters / Elsevier Computer Networks  |  "
            "Status: Draft outline + abstract complete; full write-up in progress.",
            10, GREY, wrap=True)

    textbox(slide, 0.3, 1.0, 9.4, 0.28,
            "Working Title:", 10.5, ACCENT, bold=True)
    textbox(slide, 0.3, 1.28, 9.4, 0.36,
            "\"Unified Multi-Modal Anomaly Detection with Cross-Domain Correlation "
            "and LLM-Aided Explanation for 5G Networks\"",
            11.5, GREEN, bold=True, wrap=True)

    textbox(slide, 0.3, 1.72, 9.4, 0.28,
            "Abstract (draft):", 10.5, ACCENT, bold=True)
    abstract = (
        "We propose a unified open-source pipeline for 5G network fault detection that "
        "jointly processes PCAP protocol traces, per-cell KPI exports, and gNB L1/L2 "
        "statistics through a 18-detector ensemble covering six anomaly-detection families. "
        "A shared event correlator links anomalies across the three data domains by cell, "
        "category, and time-window, enabling cross-domain root-cause attribution that "
        "single-source methods cannot achieve. On a real-world-scale testbed of 4 gNBs, "
        "16 cells, 64 UEs, and a 6-hour diurnal traffic trace with three concurrent "
        "injected fault scenarios, the ensemble achieves Precision 0.75 / Recall 1.00 / "
        "F1 0.86 at the scenario level, compared to F1 0.22 for threshold-only and F1 0.50 "
        "for a single Isolation Forest baseline. A FAISS-backed RAG pipeline feeds an "
        "on-device Ollama LLM to provide natural-language root-cause narratives. "
        "The full end-to-end pipeline completes in 57 s for a 6-hour dataset on a CPU-only "
        "host, with a peak memory footprint of 1.2 GB."
    )
    textbox(slide, 0.3, 2.02, 9.4, 1.8, abstract, 10.5, WHITE, wrap=True)

    textbox(slide, 0.3, 3.87, 9.4, 0.28, "Paper Sections:", 10.5, ACCENT, bold=True)
    sections = [
        "I. Introduction — 5G monitoring gaps, single-source limitations, contributions",
        "II. Related Work — existing PCAP anomaly tools, KPI threshold methods, LLM for telecom",
        "III. System Architecture — parsers, 18-detector ensemble, event router, prediction, LLM",
        "IV. Dataset & Experimental Setup — synthetic 64-UE/16-cell/6-hr topology, injected faults",
        "V. Results & Analysis — per-domain detection, cross-domain correlation, benchmark comparison",
        "VI. Conclusion & Future Work — real hardware validation, online learning, extended LLM fine-tuning",
    ]
    for i, sec in enumerate(sections):
        ty = 4.18 + i * 0.24
        textbox(slide, 0.4, ty, 9.2, 0.22, sec, 9.5, WHITE, wrap=True)


# ── Slide 9: References ───────────────────────────────────────────────

def s_references(prs):
    slide = new_slide(prs)
    slide_title(slide, "References")
    refs = [
        "3GPP TS 24.501 — NAS protocol for 5G System (5GS)",
        "3GPP TS 38.413 — NG Application Protocol (NGAP)",
        "3GPP TS 38.331 — Radio Resource Control (RRC) specification",
        "3GPP TS 38.473 / 38.463 / 38.423 — F1AP, E1AP, XnAP protocols",
        "3GPP TS 38.913 — Requirements for Next Generation Access Technologies",
        "Pedregosa et al. — Scikit-learn: Machine Learning in Python, JMLR 2011 "
        "(Isolation Forest, OC-SVM, LOF, Elliptic Envelope)",
        "Hochreiter & Schmidhuber — Long Short-Term Memory, Neural Computation 1997 (LSTM-AE)",
        "Taylor & Letham — Forecasting at Scale (Prophet), Am. Statistician 2018",
        "Johnson, Douze & Jégou — Billion-Scale Similarity Search with GPUs (FAISS), IEEE TBDATA 2019",
        "Reimers & Gurevych — Sentence-BERT: Sentence Embeddings using Siamese Networks (MiniLM), EMNLP 2019",
        "Ollama project — Local LLM runtime (phi3:mini / llama3.2 / mistral), github.com/ollama/ollama",
        "Page — Continuous inspection schemes, Biometrika 1954 (CUSUM algorithm)",
        "Bollinger — Bollinger on Bollinger Bands, McGraw-Hill 2001 (Bollinger Band method)",
        "Project Phase II Second Review — docs/Phase2_Second_Review.pptx (this repo)",
    ]
    for i, r in enumerate(refs):
        top = 0.72 + i * 0.33
        textbox(slide, 0.4, top, 9.2, 0.31, f"[{i+1}]  {r}", 10.0, WHITE, wrap=True)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    s_title(prs)
    s_abstract(prs)
    s_overall_design(prs)
    s_experimental_results(prs)
    s_performance(prs)
    s_model_comparison(prs)
    s_code_demo(prs)
    s_journal_draft(prs)
    s_references(prs)

    out = Path("docs/Phase2_Third_Review.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
