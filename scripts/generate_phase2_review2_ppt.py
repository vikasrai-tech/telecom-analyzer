"""
Generate the Phase II — Second Review PPT — Unified Telecom Analyzer.

Maps to the "Expectation from Students" checklist for Phase II / Second
Review: Title & Abstract, Modified Algorithm Design, Contribution of the
candidate, Results obtained (intermediate), References, 80% of code (Phase II).

Picks up where docs/Phase2_First_Review.pptx left off: First Review proposed
two novelties (genuine multi-method forecasting + an LLM root-cause agent)
and had them ~40% built. This deck shows both taken to ~80%: a bounded
ReAct agent with a rule-based/RAG-grounded fallback, a 9-layer causal rule
engine, a rolling-origin + lead-time forecast evaluation harness, and the
API/dashboard wiring (shared cross-session EventRouter, /agent/root-cause,
/analyze/stats) that makes it all usable end-to-end, not just importable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x2D, 0xC6, 0x53)
ORANGE   = RGBColor(0xFF, 0x94, 0x00)
RED      = RGBColor(0xFF, 0x4D, 0x4D)
GREY     = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE   = RGBColor(0xA0, 0x60, 0xFF)
PINK     = RGBColor(0xFF, 0x4D, 0x9F)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size, color, line_h=0.32):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_title(slide, text):
    textbox(slide, 0.3, 0.12, 9.4, 0.5, text, 22, ACCENT, bold=True)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    return slide


def footer(slide, n):
    textbox(slide, 0.3, 5.28, 6.0, 0.3,
            "Unified Telecom Analyzer — Phase II Second Review", 8, GREY)
    textbox(slide, 9.3, 5.28, 0.5, 0.3, str(n), 8, GREY, align=PP_ALIGN.RIGHT)


# ── Slide 1: Title ────────────────────────────────────────────────────

def s_title(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 0.75, 9.0, 0.9,
            "Unified Telecom Analyzer", 40, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 1.75, 9.0, 0.5,
            "From Proposal to Practice: A Forecasting Benchmark and an "
            "LLM Root-Cause Agent for 5G Anomalies",
            16, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.35, 9.0, 0.4,
            "PHASE II — SECOND REVIEW",
            16, GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.8, 9.0, 0.4,
            "Both First-Review novelties taken from proposal to a working, "
            "tested, and live-demoable system",
            13, ORANGE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.5, 9.0, 0.4,
            "M.Tech Data Science  |  PES University, Bangalore — Great Learning",
            13, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.95, 9.0, 0.4,
            "Candidate: Vikas  |  Repository: telecom-analyzer",
            12, GREY, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.8, 9.0, 0.4,
            "Python · scikit-learn · PyTorch · Prophet · statsmodels · Ollama (ReAct) "
            "· FAISS · Streamlit · FastAPI",
            11, GREY, align=PP_ALIGN.CENTER)


# ── Slide 2: Title & Abstract ────────────────────────────────────────

def s_abstract(prs):
    slide = new_slide(prs)
    slide_title(slide, "Title & Abstract")
    textbox(slide, 0.3, 0.65, 9.4, 0.5,
            "Genuine Forecast Evaluation and a Guardrailed LLM Root-Cause "
            "Agent for a Multi-Domain 5G Anomaly Detection Framework",
            13.5, GREEN, bold=True, wrap=True)
    text = (
        "First Review identified two gaps in Phase I: forecasting was claimed "
        "but never quantified, and anomalies were detected in isolation with no "
        "explanation of causality. This review reports both gaps closed. A "
        "rolling-origin backtest and an anomaly lead-time evaluator now score "
        "three forecasting methods (Prophet, Holt-Winters, direct multi-horizon "
        "LSTM) on MAE/RMSE/MAPE and on whether they flag injected anomalies "
        "before onset — including an honest negative finding that abrupt "
        "step-change anomalies have no learnable precursor, contrasted with a "
        "positive demonstration on a genuine leading-trend series. A bounded "
        "ReAct agent, backed by a 9-layer 5G protocol-stack causal-rule engine "
        "and RAG-grounded 3GPP citations, now explains WHY a correlated group "
        "of cross-source anomalies occurred — with a deterministic rule-based "
        "fallback so the system stays demoable with or without a live LLM. "
        "Both novelties are wired end-to-end: two new REST endpoints "
        "(/analyze/stats, /agent/root-cause) and a shared cross-upload session "
        "router in the dashboard, not just importable modules."
    )
    textbox(slide, 0.3, 1.25, 9.4, 3.6, text, 12, WHITE, wrap=True)
    footer(slide, 2)


# ── Slide 3: Recap — where First Review left off ──────────────────────

def s_recap(prs):
    slide = new_slide(prs)
    slide_title(slide, "Recap — Status at First Review")
    bullets(slide, 0.4, 0.75, 9.2, 4.3, [
        "Phase I: COMPLETE — 18-detector ensemble (PCAP + KPI + Stats), "
        "cross-source event correlation, RAG/LLM explanation, MLOps feedback "
        "loop; reviewed and demoed prior to Phase II.",
        "First Review proposed two Phase II novelties and had them ~40% built: "
        "(1) forecasting rebuilt with quantified accuracy instead of an "
        "unverified claim, (2) an LLM agent to explain causality across "
        "correlated anomalies, not just label each one.",
        "First Review's own closing slide set the bar for this review: "
        "\"Next: Second Review — 80% code, intermediate benchmark results.\"",
        "First Review already had: LSTM + Prophet point forecasters, an early "
        "rolling-origin harness, the honest abrupt-vs-gradual lead-time "
        "finding, a first ReAct loop, and one dashboard panel — all re-used "
        "and hardened below rather than rebuilt.",
    ], 12, WHITE, line_h=0.55)
    footer(slide, 3)


# ── Slide 4: Modified Algorithm Design ───────────────────────────────

def s_modified_algo(prs):
    slide = new_slide(prs)
    slide_title(slide, "Modified Algorithm Design")
    textbox(slide, 0.3, 0.65, 9.4, 0.35,
            "What changed structurally since First Review, per component:",
            11.5, WHITE, wrap=True)

    headers = ["Component", "First Review (proposed / partial)", "Second Review (this deck)"]
    rows = [
        ("Forecasting",
         "Prophet + LSTM point forecasts; ad-hoc accuracy checks",
         "+ Holt-Winters (damped trend) as a 3rd method; formal "
         "rolling_origin_backtest() sliding N origins per (method, cell, "
         "column) for MAE/RMSE/MAPE"),
        ("Lead-time eval",
         "Single manual pass over 3 injected anomalies",
         "evaluate_anomaly_lead_time() walks backward from each ground-truth "
         "window in fixed steps + scans for false alarms away from any window"),
        ("Root-cause reasoning",
         "Single-shot LLM prompt, no tool use, no fallback",
         "Bounded ReAct loop (max 4 steps) calling 4 typed tools; any "
         "unparsable model step falls through to a fixed deterministic tool "
         "order, never aborts"),
        ("Causal knowledge",
         "Implicit in the prompt text",
         "Explicit 9-layer causal_rules.py (PHY→...→KPI) with keyword + "
         "layer-rank matching, ported from an earlier reverted attempt and "
         "adapted to EventRouter's normalized event schema"),
        ("No-LLM fallback",
         "None — agent required Ollama",
         "rule_based_root_cause(): deterministic chain ordering by time + "
         "protocol depth, with a real RAG retrieval per hop (not a hardcoded "
         "citation string) — same 3GPP-citation guarantee with or without a "
         "live model"),
        ("System wiring",
         "Agent + forecaster importable as modules; one dashboard panel",
         "/analyze/stats + /agent/root-cause REST endpoints (blocking Ollama "
         "call offloaded via asyncio.to_thread); dashboard keeps one "
         "EventRouter per session across uploads so cross-source correlation "
         "and the agent tab see everything, not just the last file"),
    ]
    col_x = [0.3, 2.1, 5.6]
    col_w = [1.7, 3.4, 3.9]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 1.08, col_w[i], 0.3, h, 10.5, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.42 + r * 0.68
        for i, val in enumerate(row):
            textbox(slide, col_x[i], top, col_w[i], 0.64, val, 8.7, WHITE, wrap=True)
    footer(slide, 4)


# ── Slide 5: Contribution of the Candidate ────────────────────────────

def s_contribution(prs):
    slide = new_slide(prs)
    slide_title(slide, "Contribution of the Candidate")
    bullets(slide, 0.4, 0.72, 9.2, 4.5, [
        "Built src/detection/forecast_eval.py from scratch: rolling-origin "
        "backtest + anomaly lead-time evaluator + run_benchmark_report() "
        "single entry point, plus the ground-truth event table for the "
        "10-hour KPI sample used to score it.",
        "Designed and implemented src/agent/ (schemas.py, tools.py, "
        "causal_rules.py, rule_based.py, react_agent.py): a bounded ReAct "
        "agent with a parse-failure guardrail, 4 individually-testable tool "
        "wrappers, a 9-layer causal-rule engine, and a deterministic "
        "RAG-grounded fallback — adapted from a prior reverted RCA attempt "
        "(commit 6851f89) to EventRouter's current event schema.",
        "Added Holt-Winters as a 3rd forecasting method and refactored "
        "threshold/severity helpers in predictor.py so forecast_eval.py, the "
        "pipeline, and the dashboard all share one implementation.",
        "Extended the REST API (src/api/main.py): new /analyze/stats "
        "endpoint and new /agent/root-cause endpoint that combines multiple "
        "prior job results into one shared EventRouter and runs the agent "
        "over cross-source correlated groups, off the event loop.",
        "Reworked the dashboard's session state to one EventRouter shared "
        "across uploads (src/dashboard/app.py), added the Root Cause Agent "
        "panel and a simplified end-user mode, and routed predicted "
        "(not-yet-happened) anomalies into that shared router so the agent "
        "can reason about them too.",
        "Wrote 24 new tests (test_root_cause_agent.py, test_forecast_eval.py, "
        "test_api_root_cause.py, test_pipeline_prediction.py) and extended "
        "test_event_router.py / test_feedback_prediction.py for the new "
        "surface area; found and fixed 2 pre-existing bugs along the way "
        "(a PDF-export crash and a freq_s divide-by-zero).",
    ], 10.8, WHITE, line_h=0.44)
    footer(slide, 5)


# ── Slide 6: Results Obtained (Intermediate) ───────────────────────────

def s_results(prs):
    slide = new_slide(prs)
    slide_title(slide, "Results Obtained (Intermediate)")

    headers = ["Metric", "Value"]
    rows = [
        ("Full automated test suite", "83 passed, 1 skipped (needs live Ollama), 0 failed"),
        ("New tests added this review", "24 (agent 13, forecast_eval 6, API root-cause 3, "
         "pipeline prediction 2)"),
        ("Regressions vs Phase I / First Review", "0"),
        ("Forecasting methods benchmarked", "3 — Prophet, Holt-Winters, direct multi-horizon LSTM"),
        ("Forecast scoring", "Rolling-origin MAE / RMSE / MAPE per (method, cell, column) "
         "+ anomaly lead-time (detected / lead-time-h / false-alarms)"),
    ]
    col_x = [0.3, 3.6]
    col_w = [3.2, 5.9]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.68, col_w[i], 0.3, h, 11.5, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 1.02 + r * 0.5
        textbox(slide, col_x[0], top, col_w[0], 0.46, row[0], 10, WHITE, wrap=True)
        textbox(slide, col_x[1], top, col_w[1], 0.46, row[1], 10, WHITE, wrap=True)

    textbox(slide, 0.3, 3.65, 9.2, 0.3,
            "Live cross-source root-cause demo (dashboard, shared session router):",
            11, ACCENT, bold=True)
    bullets(slide, 0.4, 3.98, 9.2, 1.3, [
        "Uploaded a KPI file and a Stats file for the same cell (PCI_5) in one "
        "session; the shared EventRouter correctly correlated both sources.",
        "The agent identified dl_bler (PHY layer) as the trigger, with "
        "Handover Success Rate and Cell Availability drops correctly "
        "classified as downstream KPI-level symptoms — grounded in retrieved "
        "3GPP citations (TS 38.913, TS 38.413), not co-occurrence alone.",
    ], 10, WHITE, line_h=0.36)
    footer(slide, 6)


# ── Slide 7: Honest Forecastability Finding (carried + reconfirmed) ────

def s_finding(prs):
    slide = new_slide(prs)
    slide_title(slide, "Forecastability Finding — Reconfirmed at Second Review")
    bullets(slide, 0.4, 0.75, 9.2, 4.2, [
        "Re-run against the same 3 injected anomalies (PCI_5 congestion, "
        "PCI_8 outage, PCI_2 gradual drift) with the now-3-method ensemble "
        "(Prophet, Holt-Winters, LSTM) and the formalized "
        "evaluate_anomaly_lead_time() harness — finding holds.",
        "Abrupt step-change anomalies (congestion, outage) remain "
        "undetectable ahead of onset by pure historical extrapolation: no "
        "precursor signal exists in the data before the event starts. This is "
        "a property of the dataset (designed to validate detectors, not "
        "forecasters), not a defect in any of the 3 methods.",
        "Positive control unchanged and still passing: a synthetic series "
        "with a genuine leading trend (steady PRB-utilization ramp) is "
        "correctly flagged with real lead time — proving the harness itself "
        "works, and isolating the negative finding to the dataset's event "
        "shape rather than a bug.",
        "Why this matters for the review: an honest, reproducible boundary "
        "condition is more defensible than an unqualified forecasting claim "
        "— and it directly motivates Third Review's benchmark-vs-existing-"
        "systems comparison, where this framework's method mix can be "
        "compared against single-method baselines on the same ground truth.",
    ], 11.5, WHITE, line_h=0.5)
    footer(slide, 7)


# ── Slide 8: Disclaimer — Capabilities & Limitations ───────────────────

def s_disclaimer(prs):
    slide = new_slide(prs)
    slide_title(slide, "Disclaimer — What This System Can and Cannot Do")
    textbox(slide, 0.3, 0.62, 9.4, 0.3,
            "Scoped honestly against what has actually been run and measured "
            "as of this review — not an aspirational claim.",
            10.5, GREY, italic=True, wrap=True)

    textbox(slide, 0.3, 1.0, 4.5, 0.3, "CAN DO (validated, this review)", 12, GREEN, bold=True)
    bullets(slide, 0.3, 1.34, 4.5, 3.4, [
        "Detect anomalies across PCAP + KPI + Stats with the 18-detector "
        "ensemble (83 of 84 automated tests passing, 1 skipped, 0 failed).",
        "Correlate anomalies across sources on the same cell via the shared "
        "EventRouter — live-demoed: PCI_5 KPI+Stats correlation.",
        "Forecast KPI/Stats metrics up to 4h ahead with 3 scored methods "
        "(Prophet, Holt-Winters, LSTM), each backtested for MAE/RMSE/MAPE.",
        "Explain a correlated anomaly group's likely root cause via the "
        "ReAct LLM agent, grounded in real retrieved 3GPP citations (not "
        "co-occurrence guessing), with a deterministic rule-based fallback "
        "when no LLM is running.",
        "Serve every result through both a REST API and an interactive "
        "dashboard, end-to-end, on the generated synthetic 5G datasets.",
    ], 9.7, WHITE, line_h=0.44)

    textbox(slide, 5.1, 1.0, 4.4, 0.3, "CANNOT DO / NOT YET VALIDATED", 12, RED, bold=True)
    bullets(slide, 5.1, 1.34, 4.4, 3.4, [
        "Not validated on real operator or live network traffic — every "
        "result here is on synthetic data generated by our own scripts.",
        "Cannot forecast abrupt step-change anomalies (congestion, outage) "
        "ahead of onset — no learnable precursor exists in this data (see "
        "Forecastability Finding).",
        "Root-cause output quality depends on the causal-rule/RAG knowledge "
        "base and correlated group given to it — not a substitute for "
        "expert diagnosis, and the no-LLM rule-based path is lower-"
        "confidence than the ReAct path.",
        "No quantitative benchmark yet against published/existing anomaly-"
        "detection systems — planned for Third Review.",
        "No performance/scalability profiling at production traffic volumes "
        "yet — planned for Third Review.",
        "LLM reasoning quality depends on a locally running Ollama model "
        "(phi3:mini) being available; journal write-up not yet drafted.",
    ], 9.7, WHITE, line_h=0.44)

    textbox(slide, 0.3, 4.95, 9.2, 0.35,
            "Evidence for the left column: 83 passed / 1 skipped / 0 failed "
            "(full suite), 3 forecasting methods benchmarked, live "
            "cross-source demo (PCI_5 dl_bler -> Handover Success Rate + Cell "
            "Availability, TS 38.913 / TS 38.413 citations).",
            9, GREY, wrap=True)
    footer(slide, 8)


# ── Slide 9: Code Completion Status ───────────────────────────────────

def s_code_status(prs):
    slide = new_slide(prs)
    slide_title(slide, "Phase II Code Completion — ~80% (Second Review target)")
    headers = ["Item", "Status"]
    rows = [
        ("Forecasting: Prophet + Holt-Winters + LSTM point forecasters", "Done"),
        ("Rolling-origin backtest + anomaly lead-time evaluation harness", "Done"),
        ("9-layer causal rule engine (src/agent/causal_rules.py)", "Done"),
        ("Bounded ReAct root-cause agent + parse-failure guardrail", "Done"),
        ("Deterministic RAG-grounded rule-based fallback (no-LLM path)", "Done"),
        ("/analyze/stats + /agent/root-cause REST endpoints", "Done"),
        ("Dashboard: shared session router, Root Cause Agent panel, "
         "simple mode", "Done"),
        ("24 new automated tests, 0 regressions (83 passed / 1 skipped total)", "Done"),
        ("Benchmark vs published single-method baselines (comparison study)", "Pending — Third Review"),
        ("Performance/scalability profiling at real-world scale", "Pending — Third Review"),
        ("Journal paper draft (Phase I + II)", "Pending — Third Review"),
    ]
    col_x = [0.3, 7.3]
    col_w = [6.8, 2.4]
    for i, h in enumerate(headers):
        textbox(slide, col_x[i], 0.6, col_w[i], 0.28, h, 11, ACCENT, bold=True)
    for r, row in enumerate(rows):
        top = 0.9 + r * 0.365
        color = GREEN if row[1] == "Done" else ORANGE
        textbox(slide, col_x[0], top, col_w[0], 0.35, row[0], 9.2, WHITE, wrap=True)
        textbox(slide, col_x[1], top, col_w[1], 0.35, row[1], 9.2, color, bold=True)

    textbox(slide, 0.3, 4.98, 9.2, 0.3,
            "8 of 11 items done -> ~80% of planned Phase-II Second-Review scope "
            "complete; remaining items are Third-Review deliverables.",
            9.5, GREY, italic=True, wrap=True)


# ── Slide 9: References ───────────────────────────────────────────────

def s_references(prs):
    slide = new_slide(prs)
    slide_title(slide, "References")
    refs = [
        "3GPP TS 24.501 — Non-Access-Stratum (NAS) protocol for 5G System",
        "3GPP TS 38.413 — NG Application Protocol (NGAP)",
        "3GPP TS 38.331 — Radio Resource Control (RRC) protocol specification",
        "3GPP TS 38.473 / 38.463 / 38.423 — F1AP / E1AP / XnAP",
        "3GPP TS 38.913 — Study on scenarios and requirements for next generation access technologies",
        "3GPP TS 38.321 — Medium Access Control (MAC) protocol specification",
        "Taylor & Letham — Forecasting at Scale (Prophet), 2018",
        "Holt & Winters — Exponential smoothing / damped trend forecasting (statsmodels implementation)",
        "Hochreiter & Schmidhuber — Long Short-Term Memory (LSTM), Neural Computation 1997",
        "Yao et al. — ReAct: Synergizing Reasoning and Acting in Language Models, ICLR 2023",
        "Johnson, Douze, Jegou — Billion-scale similarity search with GPUs (FAISS), 2017",
        "Reimers & Gurevych — Sentence-BERT (all-MiniLM-L6-v2), EMNLP 2019",
        "Ollama project — local LLM runtime (phi3:mini, llama3.2, mistral)",
        "Project Phase II First Review — docs/Phase2_First_Review.pptx (this repo)",
    ]
    for i, r in enumerate(refs):
        top = 0.72 + i * 0.34
        textbox(slide, 0.4, top, 9.2, 0.32, f"[{i+1}]  {r}", 10, WHITE, wrap=True)
    footer(slide, 10)


# ── Slide 11: Thank you / next steps ────────────────────────────────────

def s_close(prs):
    slide = new_slide(prs)
    textbox(slide, 0.5, 1.6, 9.0, 0.8, "Thank You", 34, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 2.35, 9.0, 0.4, "Questions & Discussion", 16, WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.15, 9.0, 0.4,
            "Next: Third Review — 100% code, benchmark vs. published "
            "single-method baselines, scalability profiling, journal paper draft",
            12.5, ORANGE, align=PP_ALIGN.CENTER)
    footer(slide, 11)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    s_title(prs)
    s_abstract(prs)
    s_recap(prs)
    s_modified_algo(prs)
    s_contribution(prs)
    s_results(prs)
    s_finding(prs)
    s_disclaimer(prs)
    s_code_status(prs)
    s_references(prs)
    s_close(prs)

    out = Path("docs/Phase2_Second_Review.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
